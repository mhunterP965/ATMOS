#!/usr/bin/env python3
"""
Generate ATMOS_A2_IDEF0_native_vNext.pptx

A2 IDEF0 decomposition (A21 -> A22 -> A23). 11.0 x 8.5 landscape, reference
frame style. Auditable mechanism-lane matrix (M1/M21/M2/M22) with per-branch
code tags and separate invisible bottom ports. >= 10 pt fonts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0,0,0); WHITE = RGBColor(0xFF,0xFF,0xFF)
prs = Presentation(); prs.slide_width = Inches(11.0); prs.slide_height = Inches(8.5)
slide = prs.slides.add_slide(prs.slide_layouts[6]); shapes = slide.shapes
def IN(v): return Inches(v)

def add_box(node, name, l, t, w, h):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = BLACK; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in ('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = node; r0.font.size = Pt(12); r0.font.bold = True
    r0.font.color.rgb = BLACK
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = name; r1.font.size = Pt(10); r1.font.color.rgb = BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb = shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb = WHITE; tb.line.fill.background()
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK
    return tb

PORT = 0.035
def add_port(x,y):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False

def _arrow(c):
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c = shapes.add_connector(MSO_CONNECTOR.ELBOW, IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.25); c.shadow.inherit = False
    if arrow: _arrow(c)
    return c

def path(pts, src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code, x, y):
    add_label(code, x-0.31, y, 0.62, 0.24, 10, bold=True, align=PP_ALIGN.CENTER,
              fill_white=True)

# ----------------------------------------------------------------- frame
bd = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb = BLACK; bd.line.width = Pt(1.25); bd.shadow.inherit = False
add_label('A2 — Decompose Generate Local Weather State',
          0.35,0.60,10.30,0.34,15,bold=True,align=PP_ALIGN.CENTER)
nb = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(9.10),IN(7.62),IN(1.45),IN(0.30))
nb.fill.background(); nb.line.color.rgb = BLACK; nb.line.width = Pt(1.0); nb.shadow.inherit = False
nt = nb.text_frame; nt.vertical_anchor = MSO_ANCHOR.MIDDLE
nt.paragraphs[0].alignment = PP_ALIGN.CENTER
nr = nt.paragraphs[0].add_run(); nr.text = 'Node: A2'; nr.font.size = Pt(11)
nr.font.bold = True; nr.font.color.rgb = BLACK

# ----------------------------------------------------------------- boxes
add_box('A21','Prepare Model Inputs and Boundary Conditions',1.30,3.55,2.25,0.90)
add_box('A22','Execute Local Micro-Weather Estimation',4.25,3.55,2.25,0.90)
add_box('A23','Derive Aviation-Relevant Weather Fields',7.20,3.55,2.35,0.90)

CY = 4.00            # chain centerline
TOPF = 3.55; BOTF = 4.45
LBX, RBX = 0.55, 10.45

# ============================ INPUT (left) ============================
path([(LBX,CY),(1.30,CY)])                         # F1 -> A21 left
add_label('F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)',
          0.40,2.92,1.78,0.55,10, fill_white=True)

# ============================ CONTROLS (top) ============================
path([(2.40,1.55),(2.40,TOPF)])                    # C21 -> A21 top
path([(5.375,1.55),(5.375,TOPF)])                  # C22 -> A22 top
path([(8.375,1.55),(8.375,TOPF)])                  # C23 -> A23 top
add_label('C21 Model configuration; resource limits; update cadence',
          1.25,1.00,2.45,0.50,10)
add_label('C22 Execution timing; numerical stability constraints; compute/power constraints',
          3.95,0.98,2.60,0.55,10)
add_label('C23 Field derivation mappings; aviation relevance rules; reporting format',
          6.95,1.00,2.70,0.55,10)

# ============================ INTERNAL FLOWS ============================
path([(3.55,CY),(4.25,CY)], src_port=True)         # A2-F1 A21->A22
path([(6.50,CY),(7.20,CY)], src_port=True)         # A2-F2 A22->A23
add_label('A2-F1 Model Inputs & Boundary Conditions (IER-04)',
          3.05,3.02,1.95,0.50,10, fill_white=True)
add_label('A2-F2 Local Micro-Weather State Estimate / Model Output',
          6.00,3.02,1.95,0.50,10, fill_white=True)

# ============================ OUTPUT (right) ============================
path([(9.55,CY),(RBX,CY)], src_port=True)          # F2 -> right boundary
add_label('F2 Local Micro-Weather State Estimate (IER-05)',
          9.00,3.00,1.62,0.50,10, fill_white=True)
add_label('To A3 Quantify Uncertainty',9.58,4.08,1.05,0.45,10, fill_white=True)

# ============================ MECHANISM-LANE MATRIX ============================
# Distinct spines (top->down): M2, M21, M1, M22.  Stubs rise to box bottom ports;
# each branch carries a code tag.  Separate invisible bottom port per endpoint.
M2y, M21y, M1y, M22y = 5.45, 5.85, 6.25, 6.65

# bottom port x by box (3 per box)
A21_M1, A21_M21, A21_M22 = 1.80, 2.45, 3.10
A22_M2, A22_M1, A22_M22 = 4.70, 5.40, 6.05
A23_M1, A23_M21, A23_M22 = 7.65, 8.30, 8.95

# --- M2 -> A22 only ---
path([(A22_M2,M2y),(A22_M2,BOTF)]); mtag('M2', A22_M2, 4.54)

# --- M21 -> A21, A23 ---
seg(A21_M21, M21y, A23_M21, M21y)                  # M21 spine
path([(A21_M21,M21y),(A21_M21,BOTF)]); mtag('M21', A21_M21, 4.54)
path([(A23_M21,M21y),(A23_M21,BOTF)]); mtag('M21', A23_M21, 4.54)

# --- M1 -> A21, A22, A23 ---
seg(A21_M1, M1y, A23_M1, M1y)                       # M1 spine
path([(A21_M1,M1y),(A21_M1,BOTF)]); mtag('M1', A21_M1, 4.54)
path([(A22_M1,M1y),(A22_M1,BOTF)]); mtag('M1', A22_M1, 4.54)
path([(A23_M1,M1y),(A23_M1,BOTF)]); mtag('M1', A23_M1, 4.54)

# --- M22 -> A21, A22, A23 ---
seg(A21_M22, M22y, A23_M22, M22y)                   # M22 spine
path([(A21_M22,M22y),(A21_M22,BOTF)]); mtag('M22', A21_M22, 4.54)
path([(A22_M22,M22y),(A22_M22,BOTF)]); mtag('M22', A22_M22, 4.54)
path([(A23_M22,M22y),(A23_M22,BOTF)]); mtag('M22', A23_M22, 4.54)

# Long mechanism reference labels (legend) along the bottom.
add_label('M1 Platforms hosting ATMOS node compute',0.55,7.05,2.05,0.48,10, fill_white=True)
add_label('M21 ATMOS preprocessing and post-processing software',2.72,7.05,2.10,0.48,10, fill_white=True)
add_label('M2 ABLE-LBM / reduced models',4.95,7.05,1.70,0.48,10, fill_white=True)
add_label('M22 Local compute/storage',6.75,7.05,1.85,0.48,10, fill_white=True)

prs.save('ATMOS_A2_IDEF0_native_vNext.pptx')
print('saved; shapes:', len(slide.shapes._spTree))
