#!/usr/bin/env python3
"""Build ATMOS_A0_native.pptx — native, editable IDEF0 A0 decomposition with GLUED connectors.

Every connector is glued (begin_connect / end_connect) to box connection points
(rect connection sites: 0=top, 1=right, 2=bottom, 3=left), so arrows reflow when a
box is moved. Boundary arrows are glued at the box end and free at the frame edge.
A5/A6 are a stacked pair with a small horizontal offset so glued top/bottom/side
approaches to each box stay clear of the other (no auto-routing through a box).

Model source: ATMOS_A0.yaml. Internal flows F1-F5 only. Black-and-white IDEF0.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TOP, RIGHT, BOTTOM, LEFT = 0, 1, 2, 3   # rect connection-site indices

prs = Presentation()
prs.slide_width = Inches(23)
prs.slide_height = Inches(11)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def conn(kind, x1, y1, x2, y2, begin=None, end=None):
    """Create a connector; optionally glue begin/end to (shape, idx). Arrowhead at END."""
    c = slide.shapes.add_connector(kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK
    c.line.width = Pt(1.5)
    if begin is not None:
        c.begin_connect(begin[0], begin[1])
    if end is not None:
        c.end_connect(end[0], end[1])
    ln = c.line._get_or_add_ln()
    for te in ln.findall(qn('a:tailEnd')):
        ln.remove(te)
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, name, nid, w=2.6, h=1.3):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0)
    rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    label(x + w - 0.6, y + h - 0.34, 0.5, 0.28, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect

# connection-point coordinates (for initial geometry of glued ends)
def L(b): return (b.left/914400, b.top/914400 + b.height/914400/2)
def R(b): return (b.left/914400 + b.width/914400, b.top/914400 + b.height/914400/2)
def T(b): return (b.left/914400 + b.width/914400/2, b.top/914400)
def B(b): return (b.left/914400 + b.width/914400/2, b.top/914400 + b.height/914400)

# -------- frame --------
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5),
                            Inches(20.1), Inches(9.3))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25)
fr.shadow.inherit = False
FX2, FY2 = 20.6, 9.8  # frame right, bottom

# -------- title / node / classification --------
label(6.0, 0.05, 11.0, 0.4, "A0 — Decompose Conduct Air-Focused Weather Exploitation Operations",
      size=18, bold=True, align=PP_ALIGN.CENTER)
label(17.4, 0.07, 3.2, 0.3, "Distribution Statement C", size=11, align=PP_ALIGN.RIGHT)
label(20.0, 10.55, 2.8, 0.3, "Node: A0", size=12, bold=True, align=PP_ALIGN.RIGHT)

# -------- activity boxes (A5 upper, A6 lower-right: offset stacked pair) --------
A1 = box(2.4, 5.2, "Acquire Micro-Weather Observations", "A1")
A2 = box(5.6, 5.2, "Generate Local Weather State", "A2")
A3 = box(8.8, 5.2, "Quantify Uncertainty", "A3")
A4 = box(12.0, 5.2, "Federate and Maintain Federated Weather Context", "A4")
A5 = box(15.2, 3.0, "Detect Threshold Crossings (Descriptive Support)", "A5")
A6 = box(17.2, 7.3, "Produce Mission-Tailored Weather Context", "A6")

ST, EL = MSO_CONNECTOR.STRAIGHT, MSO_CONNECTOR.ELBOW

# -------- internal flows F1-F5 (glued box->box) --------
conn(ST, *R(A1), *L(A2), begin=(A1, RIGHT), end=(A2, LEFT))   # F1
conn(ST, *R(A2), *L(A3), begin=(A2, RIGHT), end=(A3, LEFT))   # F2
conn(ST, *R(A3), *L(A4), begin=(A3, RIGHT), end=(A4, LEFT))   # F3
conn(EL, *R(A4), *L(A5), begin=(A4, RIGHT), end=(A5, LEFT))   # F4
conn(EL, *R(A4), *L(A6), begin=(A4, RIGHT), end=(A6, LEFT))   # F5
label(3.9, 4.55, 2.8, 0.5, "F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)", size=9, align=PP_ALIGN.CENTER)
label(7.1, 4.35, 2.8, 0.5, "F2 Local Micro-Weather State Estimate (IER-05)", size=9, align=PP_ALIGN.CENTER)
label(9.5, 4.55, 2.6, 0.5, "F3 Confidence Bounds & Risk Envelopes (IER-09)", size=9, align=PP_ALIGN.CENTER)
label(13.7, 2.35, 2.3, 0.5, "F4 Federated Weather Context / COWP State", size=9)
label(10.6, 7.15, 2.6, 0.5, "F5 Federated Weather Context / COWP State", size=9)

# -------- inputs (free at LEFT frame, glued to box LEFT) --------
conn(EL, 0.5, 5.5, *L(A1), end=(A1, LEFT))    # I1
conn(EL, 0.5, 6.2, *L(A1), end=(A1, LEFT))    # I2
conn(ST, 0.5, 3.65, *L(A5), end=(A5, LEFT))   # I3 -> A5
conn(ST, 0.5, 7.95, *L(A6), end=(A6, LEFT))   # I3 -> A6
label(0.6, 5.12, 1.7, 0.4, "I1 Collocated Atmospheric Observations", size=9)
label(0.6, 6.25, 1.7, 0.4, "I2 Peer Platform Weather Observations", size=9)
label(0.6, 3.30, 2.9, 0.34, "I3 Mission Weather Threshold Definitions", size=9)

# -------- controls (free at TOP frame, glued to box TOP) --------
conn(ST, 13.3, 0.5, *T(A4), end=(A4, TOP))    # C3 -> A4
conn(EL, 12.6, 0.5, *T(A4), end=(A4, TOP))    # C2 -> A4
conn(ST, 18.5, 0.5, *T(A6), end=(A6, TOP))    # C2 -> A6
conn(ST, 16.5, 0.5, *T(A5), end=(A5, TOP))    # C1 -> A5
conn(EL, 18.9, 0.5, *T(A6), end=(A6, TOP))    # C1 -> A6
label(10.3, 0.13, 2.4, 0.34, "C2 OPSEC / classification guidance", size=9, align=PP_ALIGN.CENTER)
label(12.9, 0.13, 2.6, 0.34, "C3 Network availability and DDS QoS policies", size=9, align=PP_ALIGN.CENTER)
label(15.9, 0.13, 3.0, 0.34, "C1 Mission objectives and operational constraints", size=9, align=PP_ALIGN.CENTER)

# -------- outputs (glued to box RIGHT, free at RIGHT frame) --------
conn(ST, *R(A4), FX2, R(A4)[1], begin=(A4, RIGHT))   # O1
conn(ST, *R(A5), FX2, R(A5)[1], begin=(A5, RIGHT))   # O3
conn(ST, *R(A6), FX2, R(A6)[1], begin=(A6, RIGHT))   # O2
label(20.7, 5.60, 2.2, 0.5, "O1 Fused COWP State", size=9)
label(20.7, 3.40, 2.2, 0.5, "O3 Weather State Change Notifications", size=9)
label(20.7, 7.70, 2.2, 0.5, "O2 Mission-Tailored COWP Excerpts", size=9)

# -------- mechanisms (free at BOTTOM frame, glued to box BOTTOM) --------
conn(ST, 3.7, FY2, *B(A1), end=(A1, BOTTOM))    # M1 -> A1
conn(ST, 6.9, FY2, *B(A2), end=(A2, BOTTOM))    # M1 -> A2
conn(ST, 10.1, FY2, *B(A3), end=(A3, BOTTOM))   # M1 -> A3
conn(ST, 13.3, FY2, *B(A4), end=(A4, BOTTOM))   # M1 -> A4
conn(ST, *B(A5)[:1], FY2, *B(A5), end=(A5, BOTTOM))   # M1 -> A5 (clear left of A6)
conn(ST, *B(A6)[:1], FY2, *B(A6), end=(A6, BOTTOM))   # M1 -> A6
conn(EL, 7.3, FY2, *B(A2), end=(A2, BOTTOM))    # M2 -> A2
conn(EL, 10.5, FY2, *B(A3), end=(A3, BOTTOM))   # M2 -> A3
conn(EL, 13.7, FY2, *B(A4), end=(A4, BOTTOM))   # M3 -> A4
conn(EL, 18.9, FY2, *B(A6), end=(A6, BOTTOM))   # M3 -> A6
label(2.0, 9.92, 3.2, 0.32, "M1 Platforms hosting ATMOS node compute", size=9, align=PP_ALIGN.CENTER)
label(6.0, 9.92, 3.0, 0.32, "M2 ABLE-LBM / reduced models", size=9, align=PP_ALIGN.CENTER)
label(12.6, 9.92, 3.4, 0.32, "M3 DDS middleware and communications links", size=9, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A0_native.pptx")
print("saved ATMOS_A0_native.pptx; shapes:", len(slide.shapes))
