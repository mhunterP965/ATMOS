#!/usr/bin/env python3
"""Build ATMOS_A_MINUS_0_native.pptx — letter-size (11x8.5 landscape) IDEF0 A-0
context diagram: one function box with strict IDEF0 side semantics.

Inputs (I1,I2) enter LEFT; Controls (C1,C2,C3,I3) enter TOP; Outputs (O1,O2,O3)
leave RIGHT; Mechanisms (M1,M2,M3) enter BOTTOM. Each ICOM uses its own invisible
connector-port on the correct side; every connector is an orthogonal elbow with a
real bend (no straight/diagonal/curved connectors). Box text = node + name only.
No text below 10 pt. Native shapes; no images, no groups.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
prs = Presentation()
prs.slide_width = Inches(11.0); prs.slide_height = Inches(8.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.25)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    assert len(points) >= 3, "every ICOM must be an elbow (>=2 segments)"
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        assert abs(x1 - x2) < 1e-6 or abs(y1 - y2) < 1e-6, "non-orthogonal segment"
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def port(x, y):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x - 0.02), Inches(y - 0.02), Inches(0.04), Inches(0.04))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.TOP):
    assert size >= 10, f"font<10pt: {text!r}"
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0.5))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


# frame
FL, FT, FR, FB = 0.3, 0.7, 10.7, 8.0
fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT), Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.0); fr.shadow.inherit = False

# title / marking / node
label(0.3, 0.05, 8.1, 0.27, "A-0 — Conduct Air-Focused Weather Exploitation Operations", size=12, bold=True)
label(8.5, 0.06, 2.2, 0.24, "Distribution Statement C", size=10, align=PP_ALIGN.RIGHT)
label(9.0, 7.55, 1.6, 0.25, "Node: A-0", size=10, bold=True, align=PP_ALIGN.RIGHT)

# central function box (node + name only)
BL, BR, BT, BB = 4.0, 7.0, 3.2, 5.3
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(BL), Inches(BT), Inches(BR - BL), Inches(BB - BT))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BLACK
box.line.width = Pt(1.75); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "A-0"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Conduct Air-Focused Weather Exploitation Operations"
r2.font.size = Pt(11); r2.font.bold = False; r2.font.color.rgb = BLACK; r2.font.name = "Arial"

# ============ INPUTS (left side) ============
path([(FL, 3.45), (1.6, 3.45), (1.6, 3.7), (BL, 3.7)]); port(BL, 3.7)     # I1
path([(FL, 5.05), (1.6, 5.05), (1.6, 4.8), (BL, 4.8)]); port(BL, 4.8)     # I2
label(0.35, 3.06, 3.0, 0.32, "I1 Collocated Atmospheric Observations", size=10)
label(0.35, 5.12, 3.0, 0.32, "I2 Peer Platform Weather Observations", size=10)

# ============ CONTROLS (top side) ============
path([(2.0, FT), (2.0, 1.35), (4.45, 1.35), (4.45, BT)]); port(4.45, BT)  # C1
path([(4.2, FT), (4.2, 1.75), (5.15, 1.75), (5.15, BT)]); port(5.15, BT)  # C2
path([(7.0, FT), (7.0, 2.15), (5.85, 2.15), (5.85, BT)]); port(5.85, BT)  # C3
path([(9.3, FT), (9.3, 2.55), (6.55, 2.55), (6.55, BT)]); port(6.55, BT)  # I3 (routed as control)
label(0.85, 0.36, 2.3, 0.32, "C1 Mission objectives and operational constraints", size=10, align=PP_ALIGN.CENTER)
label(3.15, 0.36, 2.3, 0.32, "C2 OPSEC / classification guidance", size=10, align=PP_ALIGN.CENTER)
label(5.85, 0.36, 2.3, 0.32, "C3 Network availability and DDS QoS policies", size=10, align=PP_ALIGN.CENTER)
label(8.30, 0.36, 2.35, 0.32, "I3 Mission Weather Threshold Definitions", size=10, align=PP_ALIGN.CENTER)

# ============ OUTPUTS (right side) ============
port(BR, 3.65); path([(BR, 3.65), (9.2, 3.65), (9.2, 3.4), (FR, 3.4)])    # O1
port(BR, 4.25); path([(BR, 4.25), (9.5, 4.25), (9.5, 4.5), (FR, 4.5)])    # O2
port(BR, 4.85); path([(BR, 4.85), (9.2, 4.85), (9.2, 5.1), (FR, 5.1)])    # O3
label(7.15, 3.28, 3.3, 0.32, "O1 Fused COWP State", size=10)
label(7.15, 3.88, 3.3, 0.32, "O2 Mission-Tailored COWP Excerpts", size=10)
label(7.15, 4.48, 3.3, 0.32, "O3 Weather State Change Notifications", size=10)

# ============ MECHANISMS (bottom side) ============
path([(3.0, FB), (3.0, 6.8), (4.6, 6.8), (4.6, BB)]); port(4.6, BB)       # M1
path([(5.2, FB), (5.2, 7.0), (5.5, 7.0), (5.5, BB)]); port(5.5, BB)       # M2
path([(8.0, FB), (8.0, 7.2), (6.4, 7.2), (6.4, BB)]); port(6.4, BB)       # M3
label(1.75, 8.06, 2.5, 0.32, "M1 Platforms hosting ATMOS node compute", size=10, align=PP_ALIGN.CENTER)
label(4.30, 8.06, 2.3, 0.32, "M2 ABLE-LBM / reduced models", size=10, align=PP_ALIGN.CENTER)
label(6.80, 8.06, 2.6, 0.32, "M3 DDS middleware and communications links", size=10, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A_MINUS_0_native.pptx")
print("saved ATMOS_A_MINUS_0_native.pptx")
