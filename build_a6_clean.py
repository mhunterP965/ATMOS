#!/usr/bin/env python3
"""Build ATMOS_A6_native_clean.pptx — native, editable IDEF0 A6 decomposition.

Authoritative model (inline from task): parent A6 'Produce Mission-Tailored
Weather Context' -> A61, A62, A63, A64 left-to-right pipeline. Internal flows
A6-F1 (A61->A62), A6-F2 (A62->A63), A6-F3 (A63->A64). F4 parent input -> A61
left. Controls: C1 -> A62 top; C2 -> A63 & A64 top. I3 NOT rendered. Output O2
produced by A64 only. Mechanisms M1 (all four), M3 (A64), M61 (A64). COWP
appears only inside the F4 / O2 data labels.

Coordinate-placed right-angle connector paths; arrowhead on the terminal segment,
exactly on each box's outside-perimeter anchor. No glue, no images, no group.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, desc, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(8); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.30, 0.45, 0.24, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect


# ---------------- frame ----------------
FL, FT, FR, FB = 0.5, 1.1, 21.5, 11.8
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT),
                            Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# ---------------- title / marking / node ----------------
label(4.5, 0.10, 15.0, 0.32, "A6 — Decompose Produce Mission-Tailored Weather Context",
      size=16, bold=True, align=PP_ALIGN.CENTER)
label(18.0, 0.12, 3.4, 0.26, "Distribution Statement C", size=11, align=PP_ALIGN.RIGHT)
label(19.5, 11.45, 1.9, 0.3, "Node: A6", size=12, bold=True, align=PP_ALIGN.RIGHT)

# ---------------- activity boxes: chain A61 -> A62 -> A63 -> A64 ----------------
BT, BB = 5.0, 6.8
MY = 5.9   # horizontal spine (input, chain flows, output)
A61 = (2.6, 5.3); A62 = (7.3, 10.0); A63 = (12.0, 14.7); A64 = (16.7, 19.4)
cx = lambda b: (b[0] + b[1]) / 2
box(A61[0], BT, A61[1]-A61[0], BB-BT, "Ingest Federated Weather Context",
    "Receive federated weather context produced by A4 for downstream tailoring.", "A61")
box(A62[0], BT, A62[1]-A62[0], BB-BT, "Apply Mission Constraints and Relevance Filtering",
    "Filter weather data using mission objectives and operational relevance criteria.", "A62")
box(A63[0], BT, A63[1]-A63[0], BB-BT, "Apply OPSEC and Access Controls",
    "Filter or redact weather context based on classification and dissemination constraints.", "A63")
box(A64[0], BT, A64[1]-A64[0], BB-BT, "Produce and Disseminate Mission-Tailored Weather Context",
    "Package, format, and distribute mission-tailored excerpts to platforms and TOC.", "A64")

# ======================= PARENT INPUT F4 -> A61 left =======================
path([(FL, MY), (A61[0], MY)])
label(0.55, 5.16, 2.0, 0.6, "F4 Federated Weather Context / COWP State", size=8)

# ======================= INTERNAL FLOWS (chain on the spine) =======================
path([(A61[1], MY), (A62[0], MY)])   # A6-F1: A61 right -> A62 left
path([(A62[1], MY), (A63[0], MY)])   # A6-F2: A62 right -> A63 left
path([(A63[1], MY), (A64[0], MY)])   # A6-F3: A63 right -> A64 left
label(5.35, 5.18, 1.9, 0.5, "A6-F1 Federated Weather Context", size=8, align=PP_ALIGN.CENTER)
label(10.05, 5.18, 1.9, 0.5, "A6-F2 Filtered Weather Context", size=8, align=PP_ALIGN.CENTER)
label(14.75, 5.18, 1.9, 0.5, "A6-F3 Authorized Weather Context", size=8, align=PP_ALIGN.CENTER)

# ======================= CONTROLS (top boundary) =======================
# C1 -> A62 top
path([(cx(A62), FT), (cx(A62), BT)])
# C2 -> A63, A64 top (bus)
seg(15.7, FT, 15.7, 2.4)                    # C2 trunk from boundary
seg(cx(A63), 2.4, cx(A64), 2.4)             # C2 bus
path([(cx(A63), 2.4), (cx(A63), BT)])       # C2 -> A63 top
path([(cx(A64), 2.4), (cx(A64), BT)])       # C2 -> A64 top
label(6.75, 0.50, 3.8, 0.30, "C1 Mission objectives and operational constraints", size=8, align=PP_ALIGN.CENTER)
label(14.1, 0.50, 3.2, 0.30, "C2 OPSEC / classification guidance", size=8, align=PP_ALIGN.CENTER)

# ======================= PARENT OUTPUT O2 =======================
path([(A64[1], MY), (FR, MY)])
label(19.5, 5.14, 2.0, 0.5, "O2 Mission-Tailored COWP Excerpts", size=8)
label(21.6, 5.62, 2.3, 0.6, "To Platforms / TOC / Authorized Consumers", size=8, bold=True)

# ======================= MECHANISMS (bottom boundary) =======================
# M1 -> A61, A62, A63, A64 (bus)
seg(10.8, FB, 10.8, 8.0)                    # M1 trunk from boundary
seg(cx(A61), 8.0, cx(A64), 8.0)            # M1 bus
path([(cx(A61), 8.0), (cx(A61), BB)])       # M1 -> A61 bottom
path([(cx(A62), 8.0), (cx(A62), BB)])       # M1 -> A62 bottom
path([(cx(A63), 8.0), (cx(A63), BB)])       # M1 -> A63 bottom
path([(cx(A64), 8.0), (cx(A64), BB)])       # M1 -> A64 bottom
# M3 -> A64 bottom
path([(17.4, FB), (17.4, BB)])
# M61 -> A64 bottom
path([(18.9, FB), (18.9, BB)])
label(9.2, 11.95, 3.2, 0.30, "M1 Platforms hosting ATMOS node compute", size=8, align=PP_ALIGN.CENTER)
label(14.4, 11.95, 3.0, 0.30, "M3 DDS middleware and communications links", size=8, align=PP_ALIGN.CENTER)
label(16.9, 12.35, 3.2, 0.30, "M61 Packaging / serialization components", size=8, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A6_native_clean.pptx")
print("saved ATMOS_A6_native_clean.pptx; shapes:", len(slide.shapes))
