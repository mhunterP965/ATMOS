#!/usr/bin/env python3
"""Build ATMOS_A5_native_clean.pptx — native, editable IDEF0 A5 decomposition.

Authoritative model (inline from task): parent A5 'Detect Threshold Crossings'
-> A51, A52, A53, A54 in a left-to-right chain. Internal flows A5-F1 (A51->A52,
control-like), A5-F2 (A52->A53), A5-F3 (A53->A54). I3 is a CONTROL into A51.
F4 parent input branches to A52 (left), A53 (top-left), A54 (top-left); it never
routes through A51 or any box. Parent output O3 produced by A54 only. COWP
appears only inside the F4 data label.

Coordinate-placed right-angle connector paths; arrowhead on the terminal segment,
exactly on each box's outside-perimeter anchor. No glue, no images, no group.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, desc, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(8); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.30, 0.45, 0.24, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect


# ---------------- frame ----------------
FL, FT, FR, FB = 0.5, 1.1, 21.5, 11.8
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT),
                            Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# ---------------- title / marking / node ----------------
label(5.5, 0.10, 13.0, 0.32, "A5 — Decompose Detect Threshold Crossings",
      size=17, bold=True, align=PP_ALIGN.CENTER)
label(18.0, 0.12, 3.4, 0.26, "Distribution Statement C", size=11, align=PP_ALIGN.RIGHT)
label(19.5, 11.45, 1.9, 0.3, "Node: A5", size=12, bold=True, align=PP_ALIGN.RIGHT)

# ---------------- activity boxes: chain A51 -> A52 -> A53 -> A54 ----------------
BT, BB = 5.0, 6.8
MY = 5.9   # box mid (output rail)
A51 = (2.1, 4.8); A52 = (7.0, 9.7); A53 = (11.9, 14.6); A54 = (16.8, 19.5)
cx = lambda b: (b[0] + b[1]) / 2
box(A51[0], BT, A51[1]-A51[0], BB-BT, "Receive Mission Weather Threshold Definitions",
    "Ingest mission-defined weather thresholds for comparison against weather state.", "A51")
box(A52[0], BT, A52[1]-A52[0], BB-BT, "Present Weather Conditions Relevant to Aviation Risk",
    "Evaluate and expose current weather conditions against operationally relevant criteria.", "A52")
box(A53[0], BT, A53[1]-A53[0], BB-BT, "Present Temporal Weather Condition Trends",
    "Identify temporal changes in weather conditions relevant to operational thresholds.", "A53")
box(A54[0], BT, A54[1]-A54[0], BB-BT, "Notify Weather State Changes",
    "Generate and disseminate descriptive notifications when weather thresholds are crossed or states change.", "A54")

CHAINY = 5.6   # internal chain-flow rail (left/right edges)
F4Y = 6.0      # F4 left-entry height (A52 left edge)
RAILY = 3.8    # F4 top rail

# ======================= PARENT INPUT F4 (branch to A52 left, A53/A54 top-left) =======================
seg(FL, RAILY, 17.0, RAILY)                              # F4 top rail from left boundary
path([(6.0, RAILY), (6.0, F4Y), (A52[0], F4Y)])          # F4 -> A52 left perimeter
path([(12.2, RAILY), (12.2, BT)])                        # F4 -> A53 top-left perimeter
path([(17.0, RAILY), (17.0, BT)])                        # F4 -> A54 top-left perimeter
label(0.55, 3.05, 2.3, 0.6, "F4 Federated Weather Context / COWP State", size=8)

# ======================= CONTROLS (top boundary) =======================
# I3 -> A51 top
path([(cx(A51), FT), (cx(A51), BT)])
# C1 -> A52, A53 top (bus)
seg(10.8, FT, 10.8, 2.4)                                 # C1 trunk from boundary
seg(cx(A52), 2.4, cx(A53), 2.4)                          # C1 bus
path([(cx(A52), 2.4), (cx(A52), BT)])                    # C1 -> A52 top
path([(cx(A53), 2.4), (cx(A53), BT)])                    # C1 -> A53 top
# C2 -> A54 top
path([(cx(A54), FT), (cx(A54), BT)])
label(1.75, 0.50, 3.4, 0.30, "I3 Mission Weather Threshold Definitions", size=8, align=PP_ALIGN.CENTER)
label(8.9, 0.50, 3.8, 0.30, "C1 Mission objectives and operational constraints", size=8, align=PP_ALIGN.CENTER)
label(16.55, 0.50, 3.2, 0.30, "C2 OPSEC / classification guidance", size=8, align=PP_ALIGN.CENTER)

# ======================= INTERNAL FLOWS =======================
# A5-F1: A51 right -> A52 top (control-like, over the gap)
path([(A51[1], 5.4), (5.4, 5.4), (5.4, 4.4), (7.3, 4.4), (7.3, BT)])
label(5.0, 3.02, 2.9, 0.40, "A5-F1 Mission Weather Threshold Definitions (control)", size=8)
# A5-F2: A52 right -> A53 left
path([(A52[1], CHAINY), (A53[0], CHAINY)])
label(9.8, 5.02, 2.0, 0.52, "A5-F2 Evaluated Weather Condition State", size=7.5, align=PP_ALIGN.CENTER)
# A5-F3: A53 right -> A54 left
path([(A53[1], CHAINY), (A54[0], CHAINY)])
label(14.7, 5.02, 2.0, 0.52, "A5-F3 Detected Weather Trends / Changes", size=7.5, align=PP_ALIGN.CENTER)

# ======================= PARENT OUTPUT O3 =======================
path([(A54[1], MY), (FR, MY)])
label(19.55, 5.14, 1.95, 0.5, "O3 Weather State Change Notifications", size=8)
label(21.6, 5.62, 2.3, 0.6, "To Platforms / TOC / Downstream A6", size=8, bold=True)

# ======================= MECHANISMS (bottom boundary) =======================
# M1 -> A51, A52, A53, A54 (bus)
seg(10.8, FB, 10.8, 8.0)                                 # M1 trunk from boundary
seg(cx(A51), 8.0, cx(A54), 8.0)                          # M1 bus
path([(cx(A51), 8.0), (cx(A51), BB)])                    # M1 -> A51 bottom
path([(cx(A52), 8.0), (cx(A52), BB)])                    # M1 -> A52 bottom
path([(cx(A53), 8.0), (cx(A53), BB)])                    # M1 -> A53 bottom
path([(cx(A54), 8.0), (cx(A54), BB)])                    # M1 -> A54 bottom
# M3 -> A54 bottom only
path([(18.7, FB), (18.7, BB)])
label(9.2, 11.95, 3.2, 0.30, "M1 Platforms hosting ATMOS node compute", size=8, align=PP_ALIGN.CENTER)
label(16.0, 11.95, 3.4, 0.30, "M3 DDS middleware and communications links", size=8, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A5_native_clean.pptx")
print("saved ATMOS_A5_native_clean.pptx; shapes:", len(slide.shapes))
