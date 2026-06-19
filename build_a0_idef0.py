#!/usr/bin/env python3
"""
Generate ATMOS_A0_IDEF0_native_vNext.pptx

Native, editable PowerPoint A0 IDEF0 decomposition diagram.
- US letter landscape 11.0 x 8.5 in
- Native rectangles, text boxes, elbow (bent) connectors, arrowheads
- Invisible connector-port shapes on box faces
- Black & white IDEF0 style, >= 10 pt fonts everywhere
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(11.0)
prs.slide_height = Inches(8.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes

# ---------------------------------------------------------------- helpers
def IN(v):
    return Inches(v)

def add_box(node, name, l, t, w=1.4, h=0.95):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = BLACK
    sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = node
    r0.font.size = Pt(12)
    r0.font.bold = True
    r0.font.color.rgb = BLACK
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = name
    r1.font.size = Pt(10)
    r1.font.bold = False
    r1.font.color.rgb = BLACK
    return sp

def add_label(text, l, t, w, h, size=10, bold=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tb = shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(1)
    tf.margin_right = Pt(1)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = BLACK
    return tb

PORT = 0.035  # invisible port square side (in)
def add_port(x, y):
    """Invisible native connector-port shape sitting on a box face."""
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          IN(x - PORT / 2), IN(y - PORT / 2),
                          IN(PORT), IN(PORT))
    sp.fill.background()       # no fill
    sp.line.fill.background()  # no line
    sp.shadow.inherit = False
    return sp

def _set_arrow(conn, end=True):
    ln = conn.line._get_or_add_ln()
    tag = 'a:tailEnd' if end else 'a:headEnd'
    el = ln.makeelement(qn(tag), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(el)

def seg(x1, y1, x2, y2, arrow=False):
    """One orthogonal elbow (bent) connector segment."""
    c = shapes.add_connector(MSO_CONNECTOR.ELBOW, IN(x1), IN(y1), IN(x2), IN(y2))
    c.line.color.rgb = BLACK
    c.line.width = Pt(1.25)
    c.shadow.inherit = False
    if arrow:
        _set_arrow(c, end=True)
    return c

def path(pts, port_at_end=True):
    """Render a polyline as elbow segments; arrowhead on final segment.
    pts: list of (x,y). Adds an invisible port at the final endpoint."""
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        last = (i == len(pts) - 2)
        seg(x1, y1, x2, y2, arrow=last)
    if port_at_end:
        add_port(*pts[-1])

# ---------------------------------------------------------------- geometry
# Function boxes (left, top); main row A1-A4, stacked A5/A6 to the right.
boxes = {
    'A1': (1.30, 3.45),
    'A2': (3.02, 3.45),
    'A3': (4.74, 3.45),
    'A4': (6.46, 3.45),
    'A5': (8.60, 2.525),
    'A6': (8.60, 4.375),
}
BW, BH = 1.4, 0.95

add_box('A1', 'Acquire Micro-Weather Observations', *boxes['A1'])
add_box('A2', 'Generate Local Weather State', *boxes['A2'])
add_box('A3', 'Quantify Uncertainty', *boxes['A3'])
add_box('A4', 'Federate and Maintain Federated Weather Context', *boxes['A4'])
add_box('A5', 'Detect Threshold Crossings (Descriptive Support)', *boxes['A5'])
add_box('A6', 'Produce Mission-Tailored Weather Context', *boxes['A6'])

# convenience edges
def L(b):  return boxes[b][0]
def R(b):  return boxes[b][0] + BW
def T(b):  return boxes[b][1]
def B(b):  return boxes[b][1] + BH
def CY(b): return boxes[b][1] + BH / 2
def CX(b): return boxes[b][0] + BW / 2

LB = 0.20    # left boundary x
RB = 0.50    # (unused)
RBND = 10.50 # right boundary x (output arrowheads)
MIDY = 3.925 # main-row center y

# ============================ INPUTS (left) ============================
# I1, I2 -> A1 left ; I3 -> A5 left & A6 left
path([(LB, 3.65), (L('A1'), 3.65)])                       # I1
path([(LB, 4.15), (L('A1'), 4.15)])                       # I2
# I3 trunk + two left-side branches
path([(LB, 3.05), (8.35, 3.05), (8.35, 3.30), (L('A5'), 3.30)])   # I3 -> A5 left
path([(8.35, 3.05), (8.35, 4.70), (L('A6'), 4.70)])              # I3 -> A6 left

add_label('I1 Collocated Atmospheric Observations', 0.04, 3.18, 1.18, 0.45, 10)
add_label('I2 Peer Platform Weather Observations', 0.04, 4.18, 1.18, 0.45, 10)
add_label('I3 Mission Weather Threshold Definitions', 0.04, 2.55, 1.30, 0.45, 10)

# ============================ CONTROLS (top) ============================
# C1 -> A5,A6 ; C2 -> A4,A5,A6 ; C3 -> A4   (top faces only)
# control bus lanes
C1y, C2y, C3y = 1.45, 1.80, 2.15

# C1 bus feed
seg(2.10, 0.95, 2.10, C1y)            # label drop to C1 bus
seg(2.10, C1y, 8.95, C1y)             # C1 bus
path([(8.95, C1y), (8.95, T('A5'))])                 # C1 -> A5 top
path([(8.20, C1y), (8.20, 4.05), (8.95, 4.05), (8.95, T('A6'))])  # C1 -> A6 top
seg(8.20, C1y, 8.95, C1y)             # short bus link to A6 drop x

# C2 bus feed
seg(5.10, 1.10, 5.10, C2y)            # label drop to C2 bus
seg(5.10, C2y, 9.55, C2y)             # C2 bus
path([(6.86, C2y), (6.86, T('A4'))])                 # C2 -> A4 top
path([(9.55, C2y), (9.55, T('A5'))])                 # C2 -> A5 top
path([(8.40, C2y), (8.40, 4.18), (9.55, 4.18), (9.55, T('A6'))])  # C2 -> A6 top
seg(8.40, C2y, 9.55, C2y)             # bus link to A6 drop x

# C3 bus feed
seg(7.46, 1.25, 7.46, C3y)            # label drop to C3 lane
path([(7.46, C3y), (7.46, T('A4'))])                 # C3 -> A4 top

add_label('C1 Mission objectives and operational constraints',
          0.60, 0.55, 3.10, 0.40, 10)
add_label('C2 OPSEC / classification guidance',
          3.85, 0.55, 2.60, 0.40, 10)
add_label('C3 Network availability and DDS QoS policies',
          6.55, 0.55, 3.25, 0.40, 10)

# ============================ INTERNAL FLOWS ============================
path([(R('A1'), MIDY), (L('A2'), MIDY)])             # F1
path([(R('A2'), MIDY), (L('A3'), MIDY)])             # F2
path([(R('A3'), MIDY), (L('A4'), MIDY)])             # F3
# F4: A4 right -> A5 left ; F5: A4 right -> A6 left
path([(R('A4'), 3.60), (8.10, 3.60), (8.10, 3.10), (L('A5'), 3.10)])  # F4
path([(R('A4'), 4.25), (8.15, 4.25), (8.15, 4.55), (L('A6'), 4.55)])  # F5

# Flow labels (10pt) placed along the connectors
add_label('F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)',
          1.55, 4.50, 1.95, 0.55, 10)
add_label('F2 Local Micro-Weather State Estimate (IER-05)',
          3.27, 4.50, 1.95, 0.55, 10)
add_label('F3 Confidence Bounds & Risk Envelopes (IER-09)',
          4.99, 4.50, 1.95, 0.55, 10)
add_label('F4 Federated Weather Context / COWP State',
          6.30, 2.62, 1.55, 0.55, 10)
add_label('F5 Federated Weather Context / COWP State',
          6.30, 5.10, 1.55, 0.55, 10)

# ============================ OUTPUTS (right) ============================
# O1 from A4 (center, through A5/A6 gap), O3 from A5, O2 from A6
path([(R('A4'), MIDY), (RBND, MIDY)])                # O1
path([(R('A5'), CY('A5')), (RBND, CY('A5'))])        # O3
path([(R('A6'), CY('A6')), (RBND, CY('A6'))])        # O2

add_label('O1 Fused COWP State', 10.02, 3.50, 0.95, 0.40, 10)
add_label('O3 Weather State Change Notifications', 10.02, 2.45, 0.95, 0.55, 10)
add_label('O2 Mission-Tailored COWP Excerpts', 10.02, 4.40, 0.95, 0.55, 10)

# ============================ MECHANISMS (bottom) ============================
M1y, M2y, M3y = 6.10, 6.45, 6.80

# horizontal mechanism buses
seg(1.80, M1y, 8.95, M1y)   # M1 bus
seg(3.92, M2y, 5.64, M2y)   # M2 bus
seg(2.30, M3y, 9.55, M3y)   # M3 bus

# M1 -> A1,A2,A3,A4 (straight up), A6 (straight up), A5 (gap detour)
path([(1.80, M1y), (1.80, B('A1'))])
path([(3.52, M1y), (3.52, B('A2'))])
path([(5.24, M1y), (5.24, B('A3'))])
path([(6.96, M1y), (6.96, B('A4'))])
path([(8.95, M1y), (8.95, B('A6'))])                            # M1 -> A6 bottom
path([(8.30, M1y), (8.30, 3.95), (8.95, 3.95), (8.95, B('A5'))])  # M1 -> A5 bottom

# M2 -> A2,A3
path([(3.92, M2y), (3.92, B('A2'))])
path([(5.64, M2y), (5.64, B('A3'))])

# M3 -> A1,A4,A6 (straight), A5 (gap detour)
path([(2.30, M3y), (2.30, B('A1'))])
path([(7.36, M3y), (7.36, B('A4'))])
path([(9.55, M3y), (9.55, B('A6'))])                            # M3 -> A6 bottom
path([(8.50, M3y), (8.50, 4.05), (9.55, 4.05), (9.55, B('A5'))])  # M3 -> A5 bottom

add_label('M1 Platforms hosting ATMOS node compute',
          1.00, 6.98, 2.90, 0.40, 10)
add_label('M2 ABLE-LBM / reduced models',
          4.10, 6.98, 2.40, 0.40, 10)
add_label('M3 DDS middleware and communications links',
          6.60, 6.98, 3.30, 0.40, 10)

# ============================ TITLE / METADATA ============================
add_label('A0 — Decompose Conduct Air-Focused Weather Exploitation Operations',
          0.20, 0.06, 8.40, 0.45, 15, bold=True)
add_label('Distribution Statement C', 8.70, 0.06, 2.10, 0.22, 10, bold=True,
          align=PP_ALIGN.RIGHT)
add_label('Node: A0', 8.70, 0.27, 2.10, 0.22, 10, bold=True, align=PP_ALIGN.RIGHT)

# Footer / title block
add_label('Node: A0     |     Distribution Statement C     |     '
          'Parent (A-0): Conduct Air-Focused Weather Exploitation Operations',
          0.20, 8.12, 10.60, 0.30, 10)

prs.save('ATMOS_A0_IDEF0_native_vNext.pptx')
print('saved ATMOS_A0_IDEF0_native_vNext.pptx')
print('shapes on slide:', len(slide.shapes._spTree))
