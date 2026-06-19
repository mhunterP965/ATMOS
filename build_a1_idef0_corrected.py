#!/usr/bin/env python3
"""
Generate ATMOS_A1_IDEF0_native_vNext.pptx

Native, editable PowerPoint A1 IDEF0 decomposition (3 boxes: A11,A12,A13).
11.0 x 8.5 in landscape. Reference frame style: thin black border, centered
title, Node:A1 cell lower-right. >= 10 pt fonts. Native elbow connectors,
arrowheads, invisible ports.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0,0,0); WHITE = RGBColor(0xFF,0xFF,0xFF)
prs = Presentation(); prs.slide_width = Inches(11.0); prs.slide_height = Inches(8.5)
slide = prs.slides.add_slide(prs.slide_layouts[6]); shapes = slide.shapes
def IN(v): return Inches(v)

def add_box(node, name, l, t, w, h):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = BLACK; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in ('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = node; r0.font.size = Pt(12); r0.font.bold = True
    r0.font.color.rgb = BLACK
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = name; r1.font.size = Pt(10); r1.font.color.rgb = BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb = shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb = WHITE; tb.line.fill.background()
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK
    return tb

PORT = 0.035
def add_port(x,y):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False

def _arrow(c):
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c = shapes.add_connector(MSO_CONNECTOR.ELBOW, IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.25); c.shadow.inherit = False
    if arrow: _arrow(c)
    return c

def path(pts, src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

# ----------------------------------------------------------------- frame
bd = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb = BLACK; bd.line.width = Pt(1.25); bd.shadow.inherit = False
add_label('A1 — Decompose Acquire Micro-Weather Observations',
          0.35,0.60,10.30,0.34,15,bold=True,align=PP_ALIGN.CENTER)
nb = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(9.10),IN(7.62),IN(1.45),IN(0.30))
nb.fill.background(); nb.line.color.rgb = BLACK; nb.line.width = Pt(1.0); nb.shadow.inherit = False
nt = nb.text_frame; nt.vertical_anchor = MSO_ANCHOR.MIDDLE
nt.paragraphs[0].alignment = PP_ALIGN.CENTER
nr = nt.paragraphs[0].add_run(); nr.text = 'Node: A1'; nr.font.size = Pt(11)
nr.font.bold = True; nr.font.color.rgb = BLACK

# ----------------------------------------------------------------- boxes
add_box('A11','Sense Collocated Atmospheric Conditions',1.75,2.35,2.20,0.85)
add_box('A12','Receive Peer Platform Weather Observations',1.75,5.00,2.20,0.85)
add_box('A13','Time-Tag, Geo-Tag, and Quality-Check Observations',6.00,3.65,2.65,0.95)

LBX, RBX = 0.55, 10.45

# ============================ INPUTS (left) ============================
path([(LBX,2.775),(1.75,2.775)])                  # I1 -> A11 left
path([(LBX,5.425),(1.75,5.425)])                  # I2 -> A12 left
add_label('I1 Collocated Atmospheric Observations',0.42,2.28,1.30,0.45,10)
add_label('I2 Peer Platform Weather Observations',0.42,4.90,1.30,0.45,10)

# ============================ CONTROLS (top) ============================
path([(2.85,1.57),(2.85,2.35)])                                   # C11 -> A11 top
path([(4.25,1.57),(4.25,4.45),(2.85,4.45),(2.85,5.00)])          # C12 -> A12 top
path([(7.32,1.57),(7.32,3.65)])                                  # C13 -> A13 top
add_label('C11 Sensor configuration; platform operating state; sampling policies',
          0.80,1.00,2.45,0.56,10)
add_label('C12 DDS subscription filters; QoS; communications availability',
          3.45,1.00,2.45,0.56,10)
add_label('C13 QC rules; time/position accuracy requirements',
          6.10,1.00,2.55,0.56,10)

# ============================ INTERNAL FLOWS ============================
path([(3.95,2.90),(5.10,2.90),(5.10,3.90),(6.00,3.90)], src_port=True)  # A1-F1
path([(3.95,5.30),(4.90,5.30),(4.90,4.35),(6.00,4.35)], src_port=True)  # A1-F2
add_label('A1-F1 Collocated Atmospheric Observations',4.42,2.44,1.85,0.40,10, fill_white=True)
add_label('A1-F2 Peer Platform Weather Observations',4.05,5.36,1.95,0.40,10, fill_white=True)

# ============================ OUTPUT (right) ============================
path([(8.65,4.125),(RBX,4.125)], src_port=True)   # F1 -> right boundary
add_label('F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)',
          8.72,3.55,1.78,0.55,10, fill_white=True)
add_label('To A2 Generate Local Weather State',8.72,4.22,1.78,0.40,10, fill_white=True)

# ============================ MECHANISMS (bottom) ============================
# Traceable mechanism-lane matrix: no single shared bus. Each mechanism is a
# distinct set of stubs, each terminating on its own invisible bottom port and
# carrying a small branch-level code tag for auditability.
My = 6.95   # stub origin near the bottom boundary (just above the long labels)

def mtag(code, x, y):
    add_label(code, x, y, 0.62, 0.24, 10, bold=True, fill_white=True)

# --- M1 -> A11, A12, A13 (branches to all three; each stub tagged "M1") ---
path([(2.45,My),(2.45,5.85)])                                 # M1 -> A12
path([(6.60,My),(6.60,4.60)])                                 # M1 -> A13
path([(1.50,My),(1.50,4.05),(2.45,4.05),(2.45,3.20)])         # M1 -> A11 (left detour)
mtag('M1', 2.02, 5.92)      # near A12 M1 stub
mtag('M1', 6.18, 4.66)      # near A13 M1 stub
mtag('M1', 2.02, 3.30)      # near A11 M1 stub

# --- M11 -> A11 only (dedicated stub, tagged "M11") ---
path([(4.05,My),(4.05,4.25),(3.25,4.25),(3.25,3.20)])         # M11 -> A11 (right detour)
mtag('M11', 2.92, 3.30)     # near A11 M11 stub

# --- M12 -> A12 only (dedicated stub, tagged "M12") ---
path([(3.25,My),(3.25,5.85)])                                 # M12 -> A12
mtag('M12', 2.92, 5.92)     # near A12 M12 stub

# --- M13 -> A13 only (dedicated stub, tagged "M13") ---
path([(7.90,My),(7.90,4.60)])                                 # M13 -> A13
mtag('M13', 7.48, 4.66)     # near A13 M13 stub

# Long mechanism reference labels along the bottom (legend), 10 pt.
add_label('M1 Platforms hosting ATMOS node compute',0.55,7.05,2.05,0.50,10, fill_white=True)
add_label('M11 Onboard/attached sensors',2.70,7.05,1.45,0.55,10, fill_white=True)
add_label('M12 DDS middleware; radios/network transport',4.22,7.05,2.05,0.50,10, fill_white=True)
add_label('M13 ATMOS edge processing; navigation solution (GPS/INS); local compute',
          6.42,7.02,2.60,0.58,10, fill_white=True)

prs.save('ATMOS_A1_IDEF0_native_vNext_corrected.pptx')
print('saved; shapes:', len(slide.shapes._spTree))
