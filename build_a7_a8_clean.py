#!/usr/bin/env python3
"""Build ATMOS_A7_A8_native_clean.pptx — two slides, native editable IDEF0.

Slide 1: A7 External Reference (Execute Mission Adaptation) — explicitly marked
         external to the ATMOS system boundary; NOT an internal A0 child. 2x2
         external-domain boxes A71/A72/A73/A74 with source-defined ICOMs and
         internal flows A7-F1/F2/F3.
Slide 2: A8 internal decomposition (Assess Outcomes and Learn) — A81/A82/A83
         left-to-right; internal flows A8-F1/F2.

Native shapes only: rectangles, editable text boxes, straight/elbow connector
segments, native arrowheads. No images, no groups, arrowheads land exactly on
outside perimeters; connectors never pass through a box interior (line-line
crossings are permitted by IDEF0).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)


def seg(sl, x1, y1, x2, y2, arrow=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(sl, points, arrow=True):
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(sl, x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(sl, x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(sl, x, y, w, h, name, desc, nid, namesz=11):
    rect = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(namesz); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(8); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(sl, x + w - 0.55, y + h - 0.30, 0.45, 0.24, nid, size=10, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect


def frame(sl, fb):
    fr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
                             Inches(21.0), Inches(fb - 1.1))
    fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25)
    fr.shadow.inherit = False


# ===================================================================
# SLIDE 1 — A7 EXTERNAL REFERENCE
# ===================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
FB = 11.4
frame(s1, FB)
label(s1, 3.5, 0.08, 17.0, 0.30, "A7 — External Reference: Execute Mission Adaptation",
      size=16, bold=True, align=PP_ALIGN.CENTER)
label(s1, 3.5, 0.40, 17.0, 0.24,
      "External C2 / mission-system activity — not an internal ATMOS A0 child decomposition",
      size=10, align=PP_ALIGN.CENTER)
label(s1, 0.6, 0.10, 3.0, 0.24, "Distribution Statement C", size=10)
label(s1, 17.6, 11.45, 3.8, 0.3, "Node: A7 External Reference", size=11, bold=True, align=PP_ALIGN.RIGHT)
# explicit external-boundary note in the clear right-center band
label(s1, 14.6, 5.72, 5.2, 0.4, "EXTERNAL TO ATMOS SYSTEM BOUNDARY",
      size=12, bold=True, align=PP_ALIGN.CENTER)

# boxes (2x2)
A71 = (3.0, 6.3, 2.6, 4.0)   # L,R,T,B
A72 = (3.0, 6.3, 8.0, 9.4)
A73 = (10.5, 13.8, 2.6, 4.0)
A74 = (10.5, 13.8, 8.0, 9.4)
box(s1, A71[0], A71[2], A71[1]-A71[0], A71[3]-A71[2], "Exchange Mission Adjustment Commands",
    "Exchange mission decisions between command systems and platforms.", "A71")
box(s1, A72[0], A72[2], A72[1]-A72[0], A72[3]-A72[2], "Publish Vertical Operations Execution Updates",
    "Disseminate aircraft execution state updates.", "A72")
box(s1, A73[0], A73[2], A73[1]-A73[0], A73[3]-A73[2], "Publish UAS Tasking Updates",
    "Disseminate UAS tasking decisions.", "A73")
box(s1, A74[0], A74[2], A74[1]-A74[0], A74[3]-A74[2], "Publish Coordination Status",
    "Share coordination and acknowledgement status updates.", "A74")

# ---- inputs (left boundary) ----
seg(s1, 0.5, 3.4, 1.6, 3.4)                                        # ext mission decisions trunk
path(s1, [(1.6, 3.4), (3.0, 3.4)])                                 # -> A71 left
path(s1, [(1.6, 3.4), (1.6, 5.6), (9.6, 5.6), (9.6, 3.3), (10.5, 3.3)])  # -> A73 left
path(s1, [(0.5, 8.6), (3.0, 8.6)])                                 # aircraft exec state -> A72 left
path(s1, [(0.5, 7.1), (9.4, 7.1), (9.4, 8.9), (10.5, 8.9)])        # coordination state -> A74 left
label(s1, 0.6, 2.60, 1.05, 0.45, "External mission decisions", size=7.5)
label(s1, 0.6, 8.1, 1.5, 0.4, "Aircraft execution state", size=7.5)
label(s1, 0.6, 6.55, 1.5, 0.4, "Coordination state", size=7.5)

# ---- controls (top boundary) ----
seg(s1, 6.9, 1.1, 6.9, 1.9); seg(s1, 4.0, 1.9, 11.7, 1.9)          # command authority bus
path(s1, [(4.0, 1.9), (4.0, 2.6)])                                 # -> A71 top
path(s1, [(11.7, 1.9), (11.7, 2.6)])                               # -> A73 top
path(s1, [(5.5, 1.1), (5.5, 2.6)])                                 # rules of engagement -> A71 top
seg(s1, 8.2, 1.1, 8.2, 7.6); seg(s1, 5.0, 7.6, 11.7, 7.6)          # reporting policy bus (down to mid-band)
path(s1, [(5.0, 7.6), (5.0, 8.0)])                                 # -> A72 top
path(s1, [(11.7, 7.6), (11.7, 8.0)])                               # -> A74 top
label(s1, 2.6, 0.74, 2.4, 0.3, "Command authority", size=7.5, align=PP_ALIGN.CENTER)
label(s1, 3.7, 1.16, 1.7, 0.3, "Rules of engagement", size=7.5, align=PP_ALIGN.CENTER)
label(s1, 8.4, 1.16, 2.2, 0.3, "Reporting policy", size=7.5)

# ---- internal flows ----
path(s1, [(6.3, 3.8), (10.5, 3.8)])                                # A7-F1 A71->A73
path(s1, [(12.15, 4.0), (12.15, 8.0)])                             # A7-F2 A73->A74
path(s1, [(6.3, 8.4), (10.5, 8.4)])                                # A7-F3 A72->A74
label(s1, 8.3, 3.43, 1.2, 0.36, "A7-F1 Mission decisions", size=7)
label(s1, 12.3, 5.7, 1.6, 0.3, "A7-F2 Tasking updates", size=7.5)
label(s1, 6.45, 8.92, 2.1, 0.32, "A7-F3 Execution updates", size=7)

# ---- outputs (right boundary) ----
path(s1, [(6.3, 3.2), (7.0, 3.2), (7.0, 5.0), (21.5, 5.0)])        # IER-19 from A71
path(s1, [(6.3, 8.8), (7.5, 8.8), (7.5, 6.6), (21.5, 6.6)])        # IER-20 from A72
path(s1, [(13.8, 3.4), (21.5, 3.4)])                              # IER-21 from A73
path(s1, [(13.8, 8.8), (21.5, 8.8)])                              # IER-22 from A74
label(s1, 18.9, 4.6, 2.55, 0.34, "Mission Adjustment Commands (IER-19)", size=7.5)
label(s1, 18.9, 6.26, 2.55, 0.34, "Vertical Ops Execution Updates (IER-20)", size=7.5)
label(s1, 18.9, 3.02, 2.55, 0.34, "UAS Tasking Updates (IER-21)", size=7.5)
label(s1, 18.9, 8.46, 2.55, 0.34, "Coordination Status Updates (IER-22)", size=7.5)

# ---- mechanisms (bottom boundary) ----
seg(s1, 8.6, FB, 8.6, 4.4); seg(s1, 4.5, 4.4, 11.7, 4.4)          # C2 systems/TOC bus (up to upper boxes)
path(s1, [(4.5, 4.4), (4.5, 4.0)])                                # -> A71 bottom
path(s1, [(11.7, 4.4), (11.7, 4.0)])                              # -> A73 bottom
path(s1, [(4.5, FB), (4.5, 9.4)])                                 # aircraft avionics -> A72 bottom
seg(s1, 7.0, FB, 7.0, 10.6); seg(s1, 5.5, 10.6, 12.15, 10.6)      # comms/DDS bus
path(s1, [(5.5, 10.6), (5.5, 9.4)])                               # -> A72 bottom
path(s1, [(12.15, 10.6), (12.15, 9.4)])                           # -> A74 bottom
label(s1, 7.0, 4.5, 2.3, 0.3, "C2 systems / TOC", size=7.5)
label(s1, 3.0, 11.45, 2.6, 0.3, "Aircraft avionics", size=7.5, align=PP_ALIGN.CENTER)
label(s1, 6.0, 11.45, 3.4, 0.3, "Communications network / DDS", size=7.5, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 2 — A8 INTERNAL DECOMPOSITION
# ===================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
FB2 = 11.4
frame(s2, FB2)
label(s2, 4.5, 0.10, 15.0, 0.32, "A8 — Decompose Assess Outcomes and Learn",
      size=17, bold=True, align=PP_ALIGN.CENTER)
label(s2, 18.0, 0.12, 3.4, 0.26, "Distribution Statement C", size=11, align=PP_ALIGN.RIGHT)
label(s2, 19.5, 11.45, 1.9, 0.3, "Node: A8", size=12, bold=True, align=PP_ALIGN.RIGHT)

BT, BB, MY = 5.2, 7.0, 6.1
B81 = (3.0, 6.4); B82 = (9.3, 12.7); B83 = (15.6, 19.0)
cx = lambda b: (b[0]+b[1])/2
box(s2, B81[0], BT, B81[1]-B81[0], BB-BT, "Compare Predicted vs Observed Conditions",
    "Assess accuracy of weather predictions against observations.", "A81")
box(s2, B82[0], BT, B82[1]-B82[0], BB-BT, "Update Models and Parameters",
    "Refine models based on lessons learned.", "A82")
box(s2, B83[0], BT, B83[1]-B83[0], BB-BT, "Archive Lessons Learned",
    "Store assessment results and updates for future use.", "A83")

# ---- inputs ----
path(s2, [(0.5, 5.8), (B81[0], 5.8)])                              # Fused COWP State -> A81 left
path(s2, [(0.5, 6.4), (B81[0], 6.4)])                              # Observations -> A81 left
path(s2, [(0.5, 4.4), (7.85, 4.4), (7.85, 6.4), (B82[0], 6.4)])    # Archived Updates -> A82 left
label(s2, 0.55, 5.36, 2.3, 0.4, "Fused COWP State (IER-12)", size=8)
label(s2, 0.55, 6.46, 2.3, 0.4, "Observations (IER-01 / IER-02)", size=8)
label(s2, 0.55, 3.95, 2.6, 0.4, "Archived Updates (IER-24)", size=8)

# ---- controls ----
path(s2, [(cx(B81), 1.1), (cx(B81), BT)])
path(s2, [(cx(B82), 1.1), (cx(B82), BT)])
path(s2, [(cx(B83), 1.1), (cx(B83), BT)])
label(s2, 3.0, 0.55, 3.4, 0.3, "Assessment criteria", size=8, align=PP_ALIGN.CENTER)
label(s2, 9.3, 0.55, 3.4, 0.3, "Model governance", size=8, align=PP_ALIGN.CENTER)
label(s2, 15.6, 0.55, 3.4, 0.3, "Data retention / OPSEC", size=8, align=PP_ALIGN.CENTER)

# ---- internal flows ----
path(s2, [(B81[1], 5.8), (B82[0], 5.8)])                           # A8-F1 A81->A82
path(s2, [(B82[1], MY), (B83[0], MY)])                             # A8-F2 A82->A83
label(s2, 6.6, 6.50, 2.6, 0.3, "A8-F1 Assessment Outputs", size=8, align=PP_ALIGN.CENTER)
label(s2, 12.8, 5.62, 2.7, 0.3, "A8-F2 Updated Models", size=8, align=PP_ALIGN.CENTER)

# ---- outputs ----
path(s2, [(B81[1], 5.4), (7.2, 5.4), (7.2, 4.0), (21.5, 4.0)])     # IER-23 from A81 (upper-right) -> boundary
path(s2, [(B83[1], MY), (21.5, MY)])                               # IER-24 from A83 -> boundary
label(s2, 17.6, 3.60, 3.6, 0.34, "Predicted vs Observed Assessment (IER-23)", size=8)
label(s2, 19.1, 5.40, 2.4, 0.64, "Archived Lessons Learned & Model Updates (IER-24)", size=8, bold=True)

# ---- mechanisms ----
path(s2, [(cx(B81), FB2), (cx(B81), BB)])
path(s2, [(cx(B82), FB2), (cx(B82), BB)])
path(s2, [(cx(B83), FB2), (cx(B83), BB)])
label(s2, 3.0, 11.45, 3.4, 0.3, "ATMOS analytics", size=8, align=PP_ALIGN.CENTER)
label(s2, 9.3, 11.45, 3.4, 0.3, "Model management tools", size=8, align=PP_ALIGN.CENTER)
label(s2, 15.6, 11.45, 3.4, 0.3, "TruWeather Data Warehouse", size=8, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A7_A8_native_clean.pptx")
print("saved ATMOS_A7_A8_native_clean.pptx; slides:", len(prs.slides.__iter__.__self__._sldIdLst))
