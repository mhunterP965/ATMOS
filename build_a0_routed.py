#!/usr/bin/env python3
"""Build ATMOS_A0_native_routed.pptx — native, editable IDEF0 A0 decomposition diagram.

Native PowerPoint shapes only: rectangles (activity boxes), straight axis-aligned
connector segments forming right-angle paths (arrowhead on the final segment),
and individual editable text boxes. No images, no SVG, no grouping.

Model source: ATMOS_A0.yaml. Internal flows F1-F5 only. Strict black-and-white IDEF0.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(22)
prs.slide_height = Inches(11.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def seg(x1, y1, x2, y2, arrow=False):
    """One straight, axis-aligned connector segment (optional arrowhead at end)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK
    c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        te = ln.makeelement(qn('a:tailEnd'),
                            {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(te)
    return c


def path(points, arrow=True):
    """Draw a right-angle path as straight segments; arrowhead on last segment."""
    for i in range(len(points) - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0)
    rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    # node id, lower-right corner (separate editable text box)
    label(x + w - 0.6, y + h - 0.34, 0.5, 0.28, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)


# -------- diagram boundary frame --------
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5),
                            Inches(19.0), Inches(10.25))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25)
fr.shadow.inherit = False

# -------- title / node / classification --------
label(6.0, 0.05, 10.0, 0.4, "A0 — Decompose Conduct Air-Focused Weather Exploitation Operations",
      size=18, bold=True, align=PP_ALIGN.CENTER)
label(16.6, 0.06, 2.8, 0.3, "Distribution Statement C", size=11, align=PP_ALIGN.RIGHT)
label(19.0, 11.05, 2.8, 0.3, "Node: A0", size=12, bold=True, align=PP_ALIGN.RIGHT)

# -------- activity boxes --------
box(2.4, 5.2, 2.6, 1.3, "Acquire Micro-Weather Observations", "A1")
box(5.6, 5.2, 2.6, 1.3, "Generate Local Weather State", "A2")
box(8.8, 5.2, 2.6, 1.3, "Quantify Uncertainty", "A3")
box(12.0, 5.2, 2.6, 1.3, "Federate and Maintain Federated Weather Context", "A4")
box(15.8, 3.0, 2.6, 1.3, "Detect Threshold Crossings (Descriptive Support)", "A5")
box(15.8, 7.3, 2.6, 1.3, "Produce Mission-Tailored Weather Context", "A6")

# -------- internal flows F1-F5 --------
path([(5.0, 5.85), (5.6, 5.85)])
path([(8.2, 5.85), (8.8, 5.85)])
path([(11.4, 5.85), (12.0, 5.85)])
path([(14.6, 5.6), (15.2, 5.6), (15.2, 3.65), (15.8, 3.65)])   # F4 -> A5
path([(14.6, 6.1), (15.2, 6.1), (15.2, 7.95), (15.8, 7.95)])   # F5 -> A6
label(3.9, 4.55, 2.8, 0.5, "F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)", size=9, align=PP_ALIGN.CENTER)
label(7.1, 4.35, 2.8, 0.5, "F2 Local Micro-Weather State Estimate (IER-05)", size=9, align=PP_ALIGN.CENTER)
label(10.3, 4.55, 2.8, 0.5, "F3 Confidence Bounds & Risk Envelopes (IER-09)", size=9, align=PP_ALIGN.CENTER)
label(14.5, 2.5, 2.3, 0.5, "F4 Federated Weather Context / COWP State", size=9)
label(15.8, 6.82, 2.4, 0.45, "F5 Federated Weather Context / COWP State", size=9)

# -------- inputs (from LEFT boundary) --------
path([(0.5, 5.6), (2.4, 5.6)])     # I1 -> A1
path([(0.5, 6.1), (2.4, 6.1)])     # I2 -> A1
path([(0.5, 9.8), (15.3, 9.8)], arrow=False)          # I3 trunk
path([(15.0, 9.8), (15.0, 3.85), (15.8, 3.85)])        # I3 -> A5
path([(15.3, 9.8), (15.3, 8.1), (15.8, 8.1)])          # I3 -> A6
label(0.55, 5.18, 1.85, 0.40, "I1 Collocated Atmospheric Observations", size=9)
label(0.55, 6.12, 1.85, 0.40, "I2 Peer Platform Weather Observations", size=9)
label(0.55, 9.42, 2.9, 0.34, "I3 Mission Weather Threshold Definitions", size=9)

# -------- controls (from TOP boundary) --------
path([(13.3, 0.5), (13.3, 5.2)])                       # C3 -> A4
path([(12.6, 0.5), (12.6, 1.0)], arrow=False)          # C2 trunk
path([(12.6, 1.0), (12.6, 5.2)])                       # C2 -> A4
path([(12.6, 1.0), (19.0, 1.0), (19.0, 5.8), (17.4, 5.8), (17.4, 7.3)])  # C2 -> A6
path([(17.1, 0.5), (17.1, 0.9)], arrow=False)          # C1 trunk
path([(17.1, 0.9), (17.1, 3.0)])                       # C1 -> A5
path([(17.1, 0.9), (19.3, 0.9), (19.3, 6.2), (17.8, 6.2), (17.8, 7.3)])  # C1 -> A6
label(10.1, 0.15, 2.3, 0.6, "C2 OPSEC / classification guidance", size=9, align=PP_ALIGN.CENTER)
label(13.7, 0.15, 2.3, 0.6, "C3 Network availability and DDS QoS policies", size=9, align=PP_ALIGN.CENTER)
label(17.5, 0.15, 2.3, 0.6, "C1 Mission objectives and operational constraints", size=9, align=PP_ALIGN.CENTER)

# -------- outputs (to RIGHT boundary) --------
path([(14.6, 5.4), (14.8, 5.4), (14.8, 2.0), (19.5, 2.0)])  # O1 from A4
path([(18.4, 3.65), (19.5, 3.65)])                          # O3 from A5
path([(18.4, 7.95), (19.5, 7.95)])                          # O2 from A6
label(19.65, 1.78, 2.25, 0.5, "O1 Fused COWP State", size=9)
label(19.65, 3.42, 2.25, 0.5, "O3 Weather State Change Notifications", size=9)
label(19.65, 7.72, 2.25, 0.5, "O2 Mission-Tailored COWP Excerpts", size=9)

# -------- mechanisms (from BOTTOM boundary) --------
# M1 -> A1..A6
path([(3.0, 10.75), (3.0, 9.2)], arrow=False)          # M1 stub from boundary
path([(2.8, 9.2), (17.1, 9.2)], arrow=False)           # M1 bus
path([(3.0, 9.2), (3.0, 6.5)])                         # M1 -> A1
path([(6.4, 9.2), (6.4, 6.5)])                         # M1 -> A2
path([(9.6, 9.2), (9.6, 6.5)])                         # M1 -> A3
path([(12.8, 9.2), (12.8, 6.5)])                       # M1 -> A4
path([(16.6, 9.2), (16.6, 8.6)])                       # M1 -> A6
path([(15.4, 9.2), (15.4, 4.9), (17.1, 4.9), (17.1, 4.3)])  # M1 -> A5 (via gap)
# M2 -> A2, A3
path([(7.0, 10.75), (7.0, 9.6)], arrow=False)          # M2 stub
path([(6.8, 9.6), (9.2, 9.6)], arrow=False)            # M2 bus
path([(6.8, 9.6), (6.8, 6.5)])                         # M2 -> A2
path([(9.2, 9.6), (9.2, 6.5)])                         # M2 -> A3
# M3 -> A4, A6
path([(14.0, 10.75), (14.0, 10.0)], arrow=False)       # M3 stub
path([(13.6, 10.0), (17.0, 10.0)], arrow=False)        # M3 bus
path([(13.6, 10.0), (13.6, 6.5)])                      # M3 -> A4
path([(17.0, 10.0), (17.0, 8.6)])                      # M3 -> A6
label(1.8, 10.82, 3.0, 0.32, "M1 Platforms hosting ATMOS node compute", size=9, align=PP_ALIGN.CENTER)
label(5.8, 10.82, 3.0, 0.32, "M2 ABLE-LBM / reduced models", size=9, align=PP_ALIGN.CENTER)
label(12.6, 10.82, 3.4, 0.32, "M3 DDS middleware and communications links", size=9, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A0_native_routed.pptx")
print("saved ATMOS_A0_native_routed.pptx; shapes:", len(slide.shapes))
