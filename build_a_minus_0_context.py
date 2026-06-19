#!/usr/bin/env python3
"""Build ATMOS_A_MINUS_0_native_context_clean.pptx — true IDEF0 A-0 context
diagram on letter-size (11x8.5 landscape). ONE central function box with ICOMs
arranged around it. No decomposition, no child boxes, no A1..A8.

Strict IDEF0 sides via separate invisible ports: I1/I2 LEFT, I3/C1/C2/C3 TOP
(I3 routed as a control), O1/O2/O3 RIGHT->boundary, M1/M2/M3 BOTTOM. Every arrow
is an orthogonal elbow with a real bend. Box text = node + name only. >=10 pt.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
prs = Presentation()
prs.slide_width = Inches(11.0); prs.slide_height = Inches(8.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.25)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    assert len(points) >= 3, "every ICOM must be an elbow (>=2 segments)"
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        assert abs(x1 - x2) < 1e-6 or abs(y1 - y2) < 1e-6, "non-orthogonal"
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def port(x, y):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x - 0.02), Inches(y - 0.02), Inches(0.04), Inches(0.04))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.TOP):
    assert size >= 10, f"font<10pt: {text!r}"
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0.5))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


# frame
FL, FT, FR, FB = 0.3, 1.05, 10.7, 8.0
fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT), Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.0); fr.shadow.inherit = False

# header: title, subtitle, marking, node
label(0.3, 0.05, 8.1, 0.27, "ATMOS Operational Architecture Context", size=12, bold=True)
label(8.5, 0.06, 2.2, 0.24, "Distribution Statement C", size=10, align=PP_ALIGN.RIGHT)
label(0.3, 0.34, 8.8, 0.24, "A-0 — Conduct Air-Focused Weather Exploitation Operations", size=10, bold=True)
label(9.0, 7.55, 1.6, 0.25, "Node: A-0", size=10, bold=True, align=PP_ALIGN.RIGHT)

# central function box (node + name only)
BL, BR, BT, BB = 4.0, 7.0, 3.4, 5.4
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(BL), Inches(BT), Inches(BR - BL), Inches(BB - BT))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BLACK
box.line.width = Pt(1.75); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "A-0"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Conduct Air-Focused Weather Exploitation Operations"
r2.font.size = Pt(11); r2.font.bold = False; r2.font.color.rgb = BLACK; r2.font.name = "Arial"

# ===== INPUTS (left) =====
path([(FL, 3.60), (1.6, 3.60), (1.6, 3.85), (BL, 3.85)]); port(BL, 3.85)   # I1
path([(FL, 5.20), (1.6, 5.20), (1.6, 4.95), (BL, 4.95)]); port(BL, 4.95)   # I2
label(0.35, 3.20, 3.0, 0.32, "I1 Collocated Atmospheric Observations", size=10)
label(0.35, 5.27, 3.0, 0.32, "I2 Peer Platform Weather Observations", size=10)

# ===== CONTROLS (top) =====
path([(2.0, FT), (2.0, 1.75), (4.45, 1.75), (4.45, BT)]); port(4.45, BT)   # C1
path([(4.2, FT), (4.2, 2.15), (5.15, 2.15), (5.15, BT)]); port(5.15, BT)   # C2
path([(7.0, FT), (7.0, 2.55), (5.85, 2.55), (5.85, BT)]); port(5.85, BT)   # C3
path([(9.3, FT), (9.3, 2.95), (6.55, 2.95), (6.55, BT)]); port(6.55, BT)   # I3 (as control)
label(0.85, 0.60, 2.3, 0.34, "C1 Mission objectives and operational constraints", size=10, align=PP_ALIGN.CENTER)
label(3.15, 0.60, 2.3, 0.34, "C2 OPSEC / classification guidance", size=10, align=PP_ALIGN.CENTER)
label(5.85, 0.60, 2.3, 0.34, "C3 Network availability and DDS QoS policies", size=10, align=PP_ALIGN.CENTER)
label(8.30, 0.60, 2.35, 0.34, "I3 Mission Weather Threshold Definitions", size=10, align=PP_ALIGN.CENTER)

# ===== OUTPUTS (right) =====
port(BR, 3.75); path([(BR, 3.75), (9.2, 3.75), (9.2, 3.50), (FR, 3.50)])   # O1
port(BR, 4.40); path([(BR, 4.40), (9.5, 4.40), (9.5, 4.65), (FR, 4.65)])   # O2
port(BR, 5.05); path([(BR, 5.05), (9.2, 5.05), (9.2, 5.30), (FR, 5.30)])   # O3
label(7.15, 3.40, 3.3, 0.32, "O1 Fused COWP State", size=10)
label(7.15, 4.03, 3.3, 0.32, "O2 Mission-Tailored COWP Excerpts", size=10)
label(7.15, 4.68, 3.3, 0.32, "O3 Weather State Change Notifications", size=10)

# ===== MECHANISMS (bottom) =====
path([(3.0, FB), (3.0, 6.90), (4.6, 6.90), (4.6, BB)]); port(4.6, BB)      # M1
path([(5.2, FB), (5.2, 7.10), (5.5, 7.10), (5.5, BB)]); port(5.5, BB)      # M2
path([(8.0, FB), (8.0, 7.30), (6.4, 7.30), (6.4, BB)]); port(6.4, BB)      # M3
label(1.75, 8.06, 2.5, 0.32, "M1 Platforms hosting ATMOS node compute", size=10, align=PP_ALIGN.CENTER)
label(4.30, 8.06, 2.3, 0.32, "M2 ABLE-LBM / reduced models", size=10, align=PP_ALIGN.CENTER)
label(6.80, 8.06, 2.6, 0.32, "M3 DDS middleware and communications links", size=10, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A_MINUS_0_native_context_clean.pptx")
print("saved ATMOS_A_MINUS_0_native_context_clean.pptx")
