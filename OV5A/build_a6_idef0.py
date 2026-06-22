#!/usr/bin/env python3
"""
Generate ATMOS_A6_IDEF0_native_vNext.pptx
A6 IDEF0 decomposition: chain A61->A62->A63->A64. 11x8.5 landscape, reference
frame style, bare 'A6' node id. F4 -> A61 left; I3 -> A62 left (around A61);
A6-F1/A6-F2/A6-F3 chain; O2 output. Auditable mechanism matrix. >=10pt.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK=RGBColor(0,0,0); WHITE=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(11.0); prs.slide_height=Inches(8.5)
slide=prs.slides.add_slide(prs.slide_layouts[6]); shapes=slide.shapes
def IN(v): return Inches(v)

def add_box(node,name,l,t,w,h):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=WHITE
    sp.line.color.rgb=BLACK; sp.line.width=Pt(1.5); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0=tf.paragraphs[0]; p0.alignment=PP_ALIGN.CENTER
    r0=p0.add_run(); r0.text=node; r0.font.size=Pt(12); r0.font.bold=True; r0.font.color.rgb=BLACK
    p1=tf.add_paragraph(); p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text=name; r1.font.size=Pt(10); r1.font.color.rgb=BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb=shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb=WHITE; tb.line.fill.background()
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=BLACK
    return tb

PORT=0.035
def add_port(x,y):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit=False

def _arrow(c):
    ln=c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c=shapes.add_connector(MSO_CONNECTOR.ELBOW,IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb=BLACK; c.line.width=Pt(1.25); c.shadow.inherit=False
    if arrow:_arrow(c)
    return c

def path(pts,src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code,x,y):
    add_label(code,x-0.33,y,0.66,0.24,10,bold=True,align=PP_ALIGN.CENTER,fill_white=True)

# ---------------- frame ----------------
bd=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb=BLACK; bd.line.width=Pt(1.25); bd.shadow.inherit=False
add_label('A6 — Decompose Produce Mission-Tailored Weather Context',
          0.35,0.60,10.30,0.34,15,bold=True,align=PP_ALIGN.CENTER)
nb=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(9.35),IN(7.62),IN(1.20),IN(0.30))
nb.fill.background(); nb.line.color.rgb=BLACK; nb.line.width=Pt(1.0); nb.shadow.inherit=False
nt=nb.text_frame; nt.vertical_anchor=MSO_ANCHOR.MIDDLE; nt.paragraphs[0].alignment=PP_ALIGN.CENTER
nr=nt.paragraphs[0].add_run(); nr.text='A6'; nr.font.size=Pt(12); nr.font.bold=True; nr.font.color.rgb=BLACK

# ---------------- boxes ----------------
add_box('A61','Ingest Federated Weather Context',1.05,3.45,2.00,0.90)
add_box('A62','Apply Mission Constraints and Relevance Filtering',3.35,3.45,2.20,0.90)
add_box('A63','Apply OPSEC and Access Controls',5.85,3.45,2.05,0.90)
add_box('A64','Produce and Disseminate Mission-Tailored Weather Context',8.15,3.45,2.10,0.90)
CY=3.90; TOPF=3.45; BOTF=4.35; RBX=10.55

# ---------------- INPUTS (left) ----------------
path([(0.50,CY),(1.05,CY)])                                   # F4 -> A61 left
add_label('F4 Federated Weather Context / COWP State',0.40,3.03,1.60,0.40,10,fill_white=True)
# I3 -> A62 left (around A61), separate left port @3.65
path([(0.50,2.92),(3.20,2.92),(3.20,3.65),(3.35,3.65)])       # I3 -> A62 left
add_label('I3 Mission Weather Threshold Definitions',1.55,2.54,1.95,0.34,10,fill_white=True)

# ---------------- CONTROLS (top) ----------------
path([(2.05,1.55),(2.05,TOPF)])                               # C61 -> A61
path([(3.90,1.55),(3.90,TOPF)])                               # C1 -> A62
path([(4.90,1.95),(4.90,TOPF)])                               # C62 -> A62
seg(6.40,1.45,8.70,1.45)                                      # C2 bus
path([(6.40,1.45),(6.40,TOPF)])                               # C2 -> A63
path([(8.70,1.45),(8.70,TOPF)])                               # C2 -> A64
path([(7.35,1.95),(7.35,TOPF)])                               # C63 -> A63
path([(9.65,1.55),(9.65,TOPF)])                               # C64 -> A64
add_label('C61 Ingestion schema; data currency and completeness rules',0.95,0.98,1.65,0.55,10)
add_label('C1 Mission objectives and operational constraints',2.70,0.98,1.55,0.55,10)
add_label('C62 Mission relevance filters; consumer role/platform profiles',4.30,1.45,1.95,0.55,10)
add_label('C2 OPSEC / classification guidance',5.95,0.98,1.45,0.55,10)
add_label('C63 Access-control rules; classification markings; dissemination constraints',7.40,1.45,2.15,0.55,10)
add_label('C64 Output formatting; dissemination policy; DDS QoS policy',8.75,0.98,1.80,0.55,10)

# ---------------- INTERNAL FLOWS (chain) ----------------
path([(3.05,CY),(3.35,CY)],src_port=True)                     # A6-F1 A61->A62 left
path([(5.55,CY),(5.85,CY)],src_port=True)                     # A6-F2 A62->A63 left
path([(7.90,CY),(8.15,CY)],src_port=True)                     # A6-F3 A63->A64 left
add_label('A6-F1 Federated Weather Context',2.50,4.80,1.85,0.34,10,fill_white=True)
add_label('A6-F2 Filtered Weather Context',5.00,4.80,1.75,0.34,10,fill_white=True)
add_label('A6-F3 Authorized Weather Context',7.05,4.80,1.70,0.34,10,fill_white=True)

# ---------------- OUTPUT (right) ----------------
path([(10.25,CY),(RBX,CY)],src_port=True)                     # O2 -> right boundary
add_label('O2 Mission-Tailored COWP Excerpts (IER-18)',8.30,2.92,2.05,0.40,10,fill_white=True)
add_label('To Platforms / TOC / Authorized Consumers',8.85,4.80,1.75,0.40,10,fill_white=True)

# ---------------- MECHANISM-LANE MATRIX ----------------
SING=5.25; M61y=5.65; M1y=6.05
A61_M1,A61_M61=1.55,2.55
A62_M1,A62_M61=3.85,5.05
A63_M1,A63_M61=6.35,7.40
A64_M1,A64_M3,A64_M62,A64_M63=8.55,9.10,9.65,10.10
# M1 spine (all four)
seg(A61_M1,M1y,A64_M1,M1y)
for x in (A61_M1,A62_M1,A63_M1,A64_M1):
    path([(x,M1y),(x,BOTF)]); mtag('M1',x,4.44)
# M61 spine (A61,A62,A63)
seg(A61_M61,M61y,A63_M61,M61y)
for x in (A61_M61,A62_M61,A63_M61):
    path([(x,M61y),(x,BOTF)]); mtag('M61',x,4.44)
# singles -> A64
path([(A64_M3,SING),(A64_M3,BOTF)]); mtag('M3',A64_M3,4.44)
path([(A64_M62,SING),(A64_M62,BOTF)]); mtag('M62',A64_M62,4.44)
path([(A64_M63,SING),(A64_M63,BOTF)]); mtag('M63',A64_M63,4.44)
# long labels (legend)
add_label('M1 Platforms hosting ATMOS node compute',0.50,6.45,1.85,0.55,10,fill_white=True)
add_label('M3 DDS middleware and communications links',2.45,6.45,1.95,0.55,10,fill_white=True)
add_label('M61 ATMOS extraction/filtering services',4.50,6.45,1.85,0.55,10,fill_white=True)
add_label('M62 Packaging / serialization components',6.45,6.45,1.85,0.55,10,fill_white=True)
add_label('M63 TOC/GCS display systems and authorized consumer interfaces',8.40,6.42,2.05,0.95,10,fill_white=True)

prs.save('ATMOS_A6_IDEF0_native_vNext.pptx')
print('saved; shapes:',len(slide.shapes._spTree))
