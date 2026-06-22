#!/usr/bin/env python3
"""Build ATMOS_A3_native_clean.pptx — native, editable IDEF0 A3 decomposition.

Authoritative model (inline from task): parent A3 'Quantify Uncertainty'
-> A31, A32, A33. Internal flows A3-F1 (A31->A32), A3-F2 (A32->A33) only;
A3-F3 NOT drawn (combination folded into A33). A34/TBD NOT drawn. COWP NOT drawn.
Parent input F2 (IER-05) branches to A32 and A33 left perimeters WITHOUT passing
through any box (top rail dropping into the inter-box gaps). Parent output F3
(IER-09) exits from A33 to the right boundary.

Coordinate-placed right-angle connector paths; the terminal segment carries the
arrowhead exactly on each box's outside-perimeter anchor. No glue, no images,
no whole-diagram group, no connector grouped with a box.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    """Draw a right-angle polyline; arrowhead only on the final segment."""
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, desc, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(9); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.32, 0.45, 0.26, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect


# -------- frame: x 0.5..21.5, y 1.1..11.8 --------
FL, FT, FR, FB = 0.5, 1.1, 21.5, 11.8
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT),
                            Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# -------- title / marking / node --------
label(6.0, 0.12, 12.0, 0.42, "A3 — Decompose Quantify Uncertainty",
      size=18, bold=True, align=PP_ALIGN.CENTER)
label(17.9, 0.16, 3.5, 0.3, "Distribution Statement C", size=12, align=PP_ALIGN.RIGHT)
label(19.5, 11.45, 1.9, 0.3, "Node: A3", size=12, bold=True, align=PP_ALIGN.RIGHT)

# -------- activity boxes: horizontal chain A31 -> A32 -> A33 --------
BY1, BY2 = 4.2, 5.8
A31 = (2.5, 5.9); A32 = (9.0, 12.4); A33 = (15.5, 18.9)
box(A31[0], BY1, A31[1] - A31[0], BY2 - BY1,
    "Estimate Observation Uncertainty",
    "Quantify uncertainty associated with sensor observations based on sensor characteristics and conditions.", "A31")
box(A32[0], BY1, A32[1] - A32[0], BY2 - BY1,
    "Estimate Model and State Uncertainty",
    "Estimate uncertainty arising from model structure, inputs, and execution constraints.", "A32")
box(A33[0], BY1, A33[1] - A33[0], BY2 - BY1,
    "Package Confidence Bounds and Risk Envelopes",
    "Combine weather state and uncertainty into packaged confidence bounds and risk envelopes.", "A33")

FY = 5.3   # internal-flow rail (lower band of left/right edges)
IY = 4.6   # F2 parent-input entry height on A32/A33 left edges (upper band)
TY = 3.6   # F2 top rail (above box tops, below control descent)

# -------- parent input F2 (LEFT boundary, branch to A32 & A33; not through boxes) --------
seg(FL, TY, 14.3, TY)                       # top rail from left boundary (no arrow)
path([(8.3, TY), (8.3, IY), (A32[0], IY)])  # branch -> A32 left perimeter
path([(14.3, TY), (14.3, IY), (A33[0], IY)])# branch -> A33 left perimeter
label(0.55, 3.06, 2.3, 0.5,
      "F2 Local Micro-Weather State Estimate (IER-05)", size=8)

# -------- internal flows A3-F1, A3-F2 (lower band y=FY) --------
path([(A31[1], FY), (A32[0], FY)])   # A3-F1: A31 right -> A32 left
path([(A32[1], FY), (A33[0], FY)])   # A3-F2: A32 right -> A33 left
label(5.98, 4.86, 2.3, 0.42,
      "A3-F1 Observation Uncertainty Metrics (IER-07)", size=8, align=PP_ALIGN.CENTER)
label(12.45, 4.86, 2.3, 0.42,
      "A3-F2 Model / State Uncertainty Estimates (IER-08)", size=8, align=PP_ALIGN.CENTER)

# -------- parent output F3 (A33 right perimeter -> RIGHT boundary) --------
path([(A33[1], 5.0), (FR, 5.0)])
label(19.0, 4.28, 2.4, 0.6,
      "F3 Confidence Bounds & Risk Envelopes (IER-09)", size=9)
label(21.6, 4.92, 2.3, 0.7, "To A4 Federate and Maintain Federated Weather Context",
      size=9, bold=True)

# -------- controls (TOP boundary -> box top perimeter) --------
C31x, C32x, C33x = 4.2, 10.7, 17.2   # box center x
path([(C31x, FT), (C31x, BY1)])
path([(C32x, FT), (C32x, BY1)])
path([(C33x, FT), (C33x, BY1)])
label(C31x - 1.6, 0.58, 3.2, 0.48,
      "C31 Sensor uncertainty models; calibration/state-of-health policies",
      size=9, align=PP_ALIGN.CENTER)
label(C32x - 1.6, 0.58, 3.2, 0.48,
      "C32 Error propagation rules; model tuning parameters",
      size=9, align=PP_ALIGN.CENTER)
label(C33x - 1.6, 0.58, 3.2, 0.48,
      "C33 Confidence representation standards; formatting constraints",
      size=9, align=PP_ALIGN.CENTER)

# -------- mechanisms (BOTTOM boundary -> box bottom perimeter) --------
# A31 bottom (2.5..5.9): M1=3.3  M31=5.0
# A32 bottom (9.0..12.4): M1=9.6 M2=10.7 M32=11.8
# A33 bottom (15.5..18.9): M1=16.2 M33=17.9

# M1 -> A31, A32, A33  (bus y=7.0, trunk x=6.5 in the A31/A32 gap)
M1y = 7.0
seg(6.5, FB, 6.5, M1y)                 # trunk from bottom boundary
seg(3.3, M1y, 16.2, M1y)               # distribution bus
path([(3.3, M1y), (3.3, BY2)])         # riser -> A31 bottom
path([(9.6, M1y), (9.6, BY2)])         # riser -> A32 bottom
path([(16.2, M1y), (16.2, BY2)])       # riser -> A33 bottom

# single-target mechanisms: direct riser from bottom boundary to box bottom
path([(5.0, FB), (5.0, BY2)])          # M31 -> A31 bottom
path([(10.7, FB), (10.7, BY2)])        # M2  -> A32 bottom
path([(11.8, FB), (11.8, BY2)])        # M32 -> A32 bottom
path([(17.9, FB), (17.9, BY2)])        # M33 -> A33 bottom

# mechanism labels (below bottom boundary; two rows to avoid overlap)
label(2.6, 11.92, 2.8, 0.5, "M31 ATMOS uncertainty module; sensor metadata",
      size=9, align=PP_ALIGN.CENTER)
label(8.7, 11.92, 2.6, 0.4, "M2 ABLE-LBM / reduced models",
      size=9, align=PP_ALIGN.CENTER)
label(16.0, 11.92, 3.0, 0.5, "M33 ATMOS packaging/serialization; DDS data types",
      size=9, align=PP_ALIGN.CENTER)
label(5.0, 12.5, 3.0, 0.4, "M1 Platforms hosting ATMOS node compute",
      size=9, align=PP_ALIGN.CENTER)
label(10.8, 12.5, 2.8, 0.4, "M32 ABLE-LBM uncertainty routines",
      size=9, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A3_native_clean.pptx")
print("saved ATMOS_A3_native_clean.pptx; shapes:", len(slide.shapes))
