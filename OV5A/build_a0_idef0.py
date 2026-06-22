#!/usr/bin/env python3
"""
Generate ATMOS_A0_IDEF0_native_vNext.pptx  (vNext / reference-styled)

Native, editable PowerPoint A0 IDEF0 decomposition diagram.
Style model: ATMOS_A0_native.pptx (visible thin black border, centered title,
Node:A0 cell lower-right, B&W). Hard layout per spec: A1-A4 main row,
A5/A6 stacked right. 11.0 x 8.5 in landscape. >= 10 pt fonts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(11.0)
prs.slide_height = Inches(8.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes

def IN(v): return Inches(v)

# ----------------------------------------------------------------- helpers
def add_box(node, name, l, t, w, h):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = BLACK; sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left','margin_right'): setattr(tf, m, Pt(2))
    for m in ('margin_top','margin_bottom'): setattr(tf, m, Pt(1))
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = node
    r0.font.size = Pt(12); r0.font.bold = True; r0.font.color.rgb = BLACK
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = name
    r1.font.size = Pt(10); r1.font.color.rgb = BLACK
    return sp

def add_label(text, l, t, w, h, size=10, bold=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, fill_white=False):
    tb = shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb = WHITE; tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = BLACK
    return tb

PORT = 0.035
def add_port(x, y):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          IN(x - PORT/2), IN(y - PORT/2), IN(PORT), IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def _arrow(conn):
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type':'triangle','w':'med','len':'med'}))

def seg(x1, y1, x2, y2, arrow=False, width=1.25):
    c = shapes.add_connector(MSO_CONNECTOR.ELBOW, IN(x1), IN(y1), IN(x2), IN(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(width); c.shadow.inherit = False
    if arrow: _arrow(c)
    return c

def path(pts):
    for i in range(len(pts)-1):
        x1,y1 = pts[i]; x2,y2 = pts[i+1]
        seg(x1, y1, x2, y2, arrow=(i == len(pts)-2))
    add_port(*pts[-1])

# ----------------------------------------------------------------- frame
# visible IDEF0 border
bd = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0.35), IN(0.55), IN(10.30), IN(7.45))
bd.fill.background(); bd.line.color.rgb = BLACK; bd.line.width = Pt(1.25)
bd.shadow.inherit = False

add_label('A0 — Decompose Conduct Air-Focused Weather Exploitation Operations',
          0.35, 0.60, 10.30, 0.34, 15, bold=True, align=PP_ALIGN.CENTER)

# Node:A0 cell, lower-right title-block style
nb = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(9.10), IN(7.62), IN(1.45), IN(0.30))
nb.fill.background(); nb.line.color.rgb = BLACK; nb.line.width = Pt(1.0)
nb.shadow.inherit = False
ntf = nb.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
nr = ntf.paragraphs[0].add_run(); nr.text = 'Node: A0'
nr.font.size = Pt(11); nr.font.bold = True; nr.font.color.rgb = BLACK

# ----------------------------------------------------------------- boxes
B = {'A1':(1.15,3.70,1.35,0.75),'A2':(2.80,3.70,1.35,0.75),
     'A3':(4.45,3.70,1.35,0.75),'A4':(6.05,3.70,1.55,0.75),
     'A5':(8.25,2.35,1.50,0.75),'A6':(8.25,5.25,1.50,0.75)}
add_box('A1','Acquire Micro-Weather Observations',*B['A1'])
add_box('A2','Generate Local Weather State',*B['A2'])
add_box('A3','Quantify Uncertainty',*B['A3'])
add_box('A4','Federate and Maintain Federated Weather Context',*B['A4'])
add_box('A5','Detect Threshold Crossings (Descriptive Support)',*B['A5'])
add_box('A6','Produce Mission-Tailored Weather Context',*B['A6'])

def L(b): return B[b][0]
def R(b): return B[b][0]+B[b][2]
def T(b): return B[b][1]
def Bo(b): return B[b][1]+B[b][3]
MIDY = 4.075
LBX, RBX = 0.55, 10.55

# ============================ INPUTS (left) ============================
path([(LBX,3.90),(L('A1'),3.90)])                 # I1 -> A1
path([(LBX,4.28),(L('A1'),4.28)])                 # I2 -> A1
path([(LBX,2.30),(8.00,2.30),(8.00,2.90),(L('A5'),2.90)])  # I3 -> A5 left
path([(8.00,2.30),(8.00,5.80),(L('A6'),5.80)])             # I3 -> A6 left
add_label('I1 Collocated Atmospheric Observations',0.45,3.34,2.00,0.30,10)
add_label('I2 Peer Platform Weather Observations',0.42,4.55,1.05,0.40,10)
add_label('I3 Mission Weather Threshold Definitions',0.42,2.02,2.05,0.26,10)

# ============================ CONTROLS (top) ============================
C1y,C2y,C3y = 1.40,1.65,1.90
# C1 bus
seg(2.10,1.30,2.10,C1y); seg(2.10,C1y,10.45,C1y)
path([(8.70,C1y),(8.70,T('A5'))])                          # C1 -> A5 top
path([(10.45,C1y),(10.45,5.10),(8.70,5.10),(8.70,T('A6'))])# C1 -> A6 top
# C2 bus
seg(5.20,1.30,5.20,C2y); seg(5.20,C2y,10.58,C2y)
path([(6.45,C2y),(6.45,T('A4'))])                          # C2 -> A4 top
path([(9.40,C2y),(9.40,T('A5'))])                          # C2 -> A5 top
path([(10.58,C2y),(10.58,5.00),(9.40,5.00),(9.40,T('A6'))])# C2 -> A6 top
# C3 bus
seg(8.00,1.30,8.00,C3y); seg(7.15,C3y,8.00,C3y)
path([(7.15,C3y),(7.15,T('A4'))])                          # C3 -> A4 top
add_label('C1 Mission objectives and operational constraints',0.55,1.02,3.30,0.28,10)
add_label('C2 OPSEC / classification guidance',3.95,1.02,2.55,0.28,10)
add_label('C3 Network availability and DDS QoS policies',6.55,1.02,3.50,0.28,10)

# ============================ INTERNAL FLOWS ============================
path([(R('A1'),MIDY),(L('A2'),MIDY)])             # F1
path([(R('A2'),MIDY),(L('A3'),MIDY)])             # F2
path([(R('A3'),MIDY),(L('A4'),MIDY)])             # F3
path([(R('A4'),3.90),(7.90,3.90),(7.90,2.65),(L('A5'),2.65)])  # F4 -> A5 left
path([(R('A4'),4.28),(7.90,4.28),(7.90,5.55),(L('A6'),5.55)])  # F5 -> A6 left
# F1/F2/F3 labels centered directly above their flow gaps (no overlap: 1.6 < 1.65 pitch)
add_label('F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)',
          1.85,3.06,1.60,0.58,10, fill_white=True)
add_label('F2 Local Micro-Weather State Estimate (IER-05)',
          3.50,3.06,1.60,0.58,10, fill_white=True)
add_label('F3 Confidence Bounds & Risk Envelopes (IER-09)',
          5.13,3.06,1.60,0.58,10, fill_white=True)
add_label('F4 Federated Weather Context / COWP State',
          7.62,3.24,1.95,0.42,10, fill_white=True)
add_label('F5 Federated Weather Context / COWP State',
          7.62,4.58,1.95,0.42,10, fill_white=True)

# ============================ OUTPUTS (right) ============================
path([(R('A4'),MIDY),(RBX,MIDY)])                 # O1
path([(R('A5'),2.725),(RBX,2.725)])               # O3
path([(R('A6'),5.625),(RBX,5.625)])               # O2
add_label('O1 Fused COWP State',8.30,3.80,1.95,0.26,10, fill_white=True)
add_label('O3 Weather State Change Notifications',9.78,2.12,0.84,0.55,10, fill_white=True)
add_label('O2 Mission-Tailored COWP Excerpts',9.78,5.72,0.84,0.55,10, fill_white=True)

# ============================ MECHANISMS (bottom) ============================
M1y,M2y,M3y = 6.70,7.05,7.40
seg(1.55,M1y,9.95,M1y)            # M1 bus
seg(3.75,M2y,5.45,M2y)            # M2 bus
seg(2.10,M3y,10.15,M3y)           # M3 bus
# M1 -> A1..A4, A6 (straight), A5 (gap detour)
path([(1.55,M1y),(1.55,Bo('A1'))])
path([(3.15,M1y),(3.15,Bo('A2'))])
path([(4.85,M1y),(4.85,Bo('A3'))])
path([(6.45,M1y),(6.45,Bo('A4'))])
path([(8.70,M1y),(8.70,Bo('A6'))])
path([(9.95,M1y),(9.95,3.60),(8.70,3.60),(8.70,Bo('A5'))])
# M2 -> A2, A3
path([(3.75,M2y),(3.75,Bo('A2'))])
path([(5.45,M2y),(5.45,Bo('A3'))])
# M3 -> A1, A4, A6 (straight), A5 (gap detour)
path([(2.10,M3y),(2.10,Bo('A1'))])
path([(7.15,M3y),(7.15,Bo('A4'))])
path([(9.40,M3y),(9.40,Bo('A6'))])
path([(10.15,M3y),(10.15,3.75),(9.40,3.75),(9.40,Bo('A5'))])
add_label('M1 Platforms hosting ATMOS node compute',1.00,7.55,2.80,0.26,10)
add_label('M2 ABLE-LBM / reduced models',4.00,7.55,2.45,0.26,10)
add_label('M3 DDS middleware and communications links',6.50,7.55,2.45,0.26,10)

prs.save('ATMOS_A0_IDEF0_native_vNext.pptx')
print('saved; shapes:', len(slide.shapes._spTree))
