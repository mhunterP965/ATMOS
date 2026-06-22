#!/usr/bin/env python3
"""Build ATMOS_A_MINUS_0_native.pptx — native, editable IDEF0 A-0 context diagram.

Native PowerPoint shapes only:
 * one rectangle (function box)
 * straight, axis-aligned (orthogonal) connectors with arrowheads
 * individual editable text boxes for every label

Model source: ATMOS_A_MINUS_0.yaml (A-0 context diagram).
Strict black-and-white IDEF0 style. One slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ---- function box geometry (inches) ----
BX, BY, BW, BH = 4.87, 2.85, 3.60, 1.80
BL, BR, BT, BB = BX, BX + BW, BY, BY + BH       # left, right, top, bottom
BCX = BX + BW / 2                               # 6.67

def add_connector(x1, y1, x2, y2):
    """Straight axis-aligned connector with an arrowhead at the END (x2,y2)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    for te in ln.findall(qn('a:tailEnd')):
        ln.remove(te)
    te = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(te)
    return conn

def add_label(x, y, w, h, text, size=9, align=PP_ALIGN.LEFT, bold=False,
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK
    r.font.name = "Arial"
    return tb

# ---- title ----
add_label(3.5, 0.10, 6.33, 0.40, "ATMOS Operational Architecture Context",
          size=18, align=PP_ALIGN.CENTER, bold=True)

# ---- central function box ----
box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(BX), Inches(BY),
                             Inches(BW), Inches(BH))
box.fill.solid(); box.fill.fore_color.rgb = WHITE
box.line.color.rgb = BLACK; box.line.width = Pt(2.0)
box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Conduct Air-Focused Weather Exploitation Operations"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"

# ---- node label A-0 (lower-right interior corner, IDEF0 convention) ----
add_label(BR - 0.75, BB - 0.36, 0.65, 0.30, "A-0",
          size=12, align=PP_ALIGN.RIGHT, bold=True, anchor=MSO_ANCHOR.BOTTOM)

# ---- INPUTS (enter from left; arrowhead at box left edge) ----
in_y = [3.15, 3.75, 4.35]
inputs = ["I1 Collocated Atmospheric Observations",
          "I2 Peer Platform Weather Observations",
          "I3 Mission Weather Threshold Definitions"]
for y, txt in zip(in_y, inputs):
    add_connector(0.50, y, BL, y)
    add_label(0.55, y - 0.30, 3.9, 0.28, txt, size=9, align=PP_ALIGN.LEFT)

# ---- CONTROLS (enter from top; arrowhead at box top edge) ----
ctl_x = [5.45, 6.67, 7.90]
controls = ["C1 Mission objectives and operational constraints",
            "C2 OPSEC / classification guidance",
            "C3 Network availability and DDS QoS policies"]
for x, txt in zip(ctl_x, controls):
    add_connector(x, 0.95, x, BT)
    add_label(x - 0.65, 0.10, 1.30, 0.82, txt, size=8, align=PP_ALIGN.CENTER)

# ---- OUTPUTS (exit right; arrowhead at right diagram boundary) ----
out_y = [3.15, 3.75, 4.35]
outputs = ["O1 Fused COWP State",
           "O2 Mission-Tailored COWP Excerpts",
           "O3 Weather State Change Notifications"]
for y, txt in zip(out_y, outputs):
    add_connector(BR, y, 12.85, y)
    add_label(8.60, y - 0.30, 4.10, 0.28, txt, size=9, align=PP_ALIGN.LEFT)

# ---- MECHANISMS (enter from bottom; arrowhead at box bottom edge) ----
mec_x = [5.45, 6.67, 7.90]
mechanisms = ["M1 Platforms hosting ATMOS node compute",
              "M2 ABLE-LBM / reduced models",
              "M3 DDS middleware and communications links"]
for x, txt in zip(mec_x, mechanisms):
    add_connector(x, 6.55, x, BB)
    add_label(x - 0.65, 6.62, 1.30, 0.82, txt, size=8, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A_MINUS_0_native.pptx")
print("saved ATMOS_A_MINUS_0_native.pptx")
print("shapes on slide:", len(slide.shapes))
