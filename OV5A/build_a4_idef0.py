#!/usr/bin/env python3
"""
Generate ATMOS_A4_IDEF0_native_vNext.pptx
A4 IDEF0 decomposition (A41,A42 left column; A43 fusion hub center; A44 lower-right).
11x8.5 landscape, reference frame style. >=10pt. Feedback A4-F4 (A44->A43).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK=RGBColor(0,0,0); WHITE=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(11.0); prs.slide_height=Inches(8.5)
slide=prs.slides.add_slide(prs.slide_layouts[6]); shapes=slide.shapes
def IN(v): return Inches(v)

def add_box(node,name,l,t,w,h):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=WHITE
    sp.line.color.rgb=BLACK; sp.line.width=Pt(1.5); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0=tf.paragraphs[0]; p0.alignment=PP_ALIGN.CENTER
    r0=p0.add_run(); r0.text=node; r0.font.size=Pt(12); r0.font.bold=True; r0.font.color.rgb=BLACK
    p1=tf.add_paragraph(); p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text=name; r1.font.size=Pt(10); r1.font.color.rgb=BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb=shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb=WHITE; tb.line.fill.background()
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=BLACK
    return tb

PORT=0.035
def add_port(x,y):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit=False

def _arrow(c):
    ln=c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c=shapes.add_connector(MSO_CONNECTOR.ELBOW,IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb=BLACK; c.line.width=Pt(1.25); c.shadow.inherit=False
    if arrow:_arrow(c)
    return c

def path(pts,src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code,x,y):
    add_label(code,x-0.31,y,0.62,0.24,10,bold=True,align=PP_ALIGN.CENTER,fill_white=True)

# ---------------- frame ----------------
bd=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb=BLACK; bd.line.width=Pt(1.25); bd.shadow.inherit=False
add_label('A4 — Decompose Federate and Maintain Federated Weather Context',
          0.35,0.60,10.30,0.34,15,bold=True,align=PP_ALIGN.CENTER)
nb=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(9.35),IN(7.62),IN(1.20),IN(0.30))
nb.fill.background(); nb.line.color.rgb=BLACK; nb.line.width=Pt(1.0); nb.shadow.inherit=False
nt=nb.text_frame; nt.vertical_anchor=MSO_ANCHOR.MIDDLE; nt.paragraphs[0].alignment=PP_ALIGN.CENTER
nr=nt.paragraphs[0].add_run(); nr.text='A4'; nr.font.size=Pt(12); nr.font.bold=True; nr.font.color.rgb=BLACK

# ---------------- boxes ----------------
add_box('A41','Publish Local Weather State',1.45,2.25,2.20,0.85)
add_box('A42','Subscribe to Peer Weather States',1.45,5.00,2.20,0.85)
add_box('A43','Stitch and Blend Weather States',5.05,3.10,2.35,0.95)
add_box('A44','Maintain Degraded or Disconnected Operation',7.75,5.35,2.25,0.85)
RBX=10.45

# ---------------- INPUTS (left) ----------------
seg(0.50,2.675,1.10,2.675)                                    # F3 entry
path([(1.10,2.675),(1.45,2.675)])                             # F3 -> A41 left
path([(1.10,2.675),(1.10,3.30),(5.05,3.30)])                  # F3 -> A43 left (port @3.30)
path([(0.50,5.425),(1.45,5.425)])                             # Peer -> A42 left
add_label('F3 Confidence Bounds & Risk Envelopes (IER-09)',0.40,2.18,1.80,0.44,10,fill_white=True)
add_label('Peer Weather State Publications (IER-10 / IER-11)',0.40,4.92,1.80,0.50,10,fill_white=True)

# ---------------- CONTROLS (top) ----------------
# C3 bus (A41,A42,A44); C2 bus (A42,A43,A44); C41->A43; C42->A44
seg(2.55,1.55,8.85,1.55)                                      # C3 bus
path([(2.55,1.55),(2.55,2.25)])                               # C3 -> A41 top
path([(8.85,1.55),(8.85,5.35)])                               # C3 -> A44 top
path([(3.80,1.55),(3.80,4.55),(2.90,4.55),(2.90,5.00)])      # C3 -> A42 top (around A41)
seg(4.05,1.80,8.20,1.80)                                      # C2 bus
path([(5.85,1.80),(5.85,3.10)])                               # C2 -> A43 top
path([(8.20,1.80),(8.20,5.35)])                               # C2 -> A44 top
path([(4.05,1.80),(4.05,4.35),(2.20,4.35),(2.20,5.00)])      # C2 -> A42 top (around A41)
path([(6.75,1.55),(6.75,3.10)])                               # C41 -> A43 top
path([(9.50,1.55),(9.50,5.35)])                               # C42 -> A44 top
add_label('C3 Network availability and DDS QoS policies',0.55,1.00,2.95,0.34,10)
add_label('C2 OPSEC / classification guidance',3.95,1.00,2.20,0.34,10)
add_label('C41 Fusion heuristics; data age limits; prioritization rules',6.30,0.98,2.25,0.40,10)
add_label('C42 Cache expiration; reconnection behavior; storage limits',8.65,0.98,1.95,0.40,10)

# ---------------- INTERNAL FLOWS ----------------
path([(3.65,2.675),(4.30,2.675),(4.30,3.50),(5.05,3.50)],src_port=True)   # A4-F1 A41->A43
path([(3.65,5.425),(4.55,5.425),(4.55,3.70),(5.05,3.70)],src_port=True)   # A4-F2 A42->A43
path([(7.40,3.80),(7.57,3.80),(7.57,5.775),(7.75,5.775)],src_port=True)   # A4-F3 A43->A44
# A4-F4 feedback A44->A43 left
path([(10.00,5.775),(10.25,5.775),(10.25,4.55),(4.80,4.55),(4.80,3.90),(5.05,3.90)],src_port=True)
add_label('A4-F1 Local Weather State Publication (IER-10)',3.75,2.28,1.85,0.42,10,fill_white=True)
add_label('A4-F2 Peer Weather State Subscription (IER-11)',3.75,4.50,1.85,0.42,10,fill_white=True)
add_label('A4-F3 Fused Weather Context State (IER-12)',6.55,4.62,1.80,0.42,10,fill_white=True)
add_label('A4-F4 Cached Weather State (IER-13)',5.30,4.20,1.80,0.32,10,fill_white=True)

# ---------------- OUTPUT (right) ----------------
path([(7.40,3.45),(RBX,3.45)],src_port=True)                  # F4 -> right boundary
add_label('F4 Federated Weather Context / COWP State',8.05,3.02,1.95,0.40,10,fill_white=True)
add_label('To A5 Detect Threshold Crossings and A6 Produce Mission-Tailored Weather Context',
          8.05,3.58,2.05,0.66,10,fill_white=True)

# ---------------- MECHANISM-LANE MATRIX ----------------
# bottom ports
A41_M1,A41_M3=2.05,3.05            # bottom y3.10 (around A42)
A42_M1,A42_M3=2.05,3.05            # bottom y5.85
A43_M1,A43_M41=5.65,6.85           # bottom y4.05
A44_M1,A44_M42=8.20,9.40           # bottom y6.20
M1y=6.95; M3y=7.25
# M1 spine + branches (A41 via detour around A42)
seg(2.05,M1y,8.20,M1y)
path([(A42_M1,M1y),(A42_M1,5.85)]); mtag('M1',A42_M1,5.92)
path([(A43_M1,M1y),(A43_M1,4.05)]); mtag('M1',A43_M1,4.12)
path([(A44_M1,M1y),(A44_M1,6.20)]); mtag('M1',A44_M1,6.27)
path([(3.95,M1y),(3.95,4.55),(A41_M1,4.55),(A41_M1,3.10)]); mtag('M1',A41_M1,3.17)
# M3 spine + branches (A41 via detour, A42 direct)
seg(3.05,M3y,4.20,M3y)
path([(A42_M3,M3y),(A42_M3,5.85)]); mtag('M3',A42_M3,5.92)
path([(4.20,M3y),(4.20,4.30),(A41_M3,4.30),(A41_M3,3.10)]); mtag('M3',A41_M3,3.17)
# M41 -> A43 (single, upper lane so no spine cross)
path([(A43_M41,6.55),(A43_M41,4.05)]); mtag('M41',A43_M41,4.12)
# M42 -> A44 (single)
path([(A44_M42,6.55),(A44_M42,6.20)]); mtag('M42',A44_M42,6.27)
add_label('M1 Platforms hosting ATMOS node compute',0.55,7.48,2.45,0.34,10,fill_white=True)
add_label('M3 DDS middleware and communications links',3.10,7.48,2.70,0.34,10,fill_white=True)
add_label('M41 ATMOS federation logic',6.00,7.48,1.85,0.34,10,fill_white=True)
add_label('M42 Local storage / cache manager',7.95,7.48,2.10,0.34,10,fill_white=True)

prs.save('ATMOS_A4_IDEF0_native_vNext.pptx')
print('saved; shapes:',len(slide.shapes._spTree))
