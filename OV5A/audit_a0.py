#!/usr/bin/env python3
"""Programmatic self-audit of ATMOS_A0_IDEF0_native_vNext.pptx.
Reads the real saved XML: connector endpoints, font sizes, box content."""
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

EMU = 914400.0
prs = Presentation('ATMOS_A0_IDEF0_native_vNext.pptx')
slide = prs.slides[0]

def emu_in(v): return v / EMU

boxes = {}   # node -> (l,t,r,b)
ports = []   # (cx,cy)
conns = []   # (bx,by, ex,ey, has_arrow)
texts = []   # (text, [sizes])

BW, BH = 1.4, 0.95

for sp in slide.shapes:
    el = sp._element
    tag = el.tag.split('}')[-1]
    if tag == 'cxnSp':
        xfrm = el.find(qn('p:spPr')).find(qn('a:xfrm'))
        flipH = xfrm.get('flipH') == '1'
        flipV = xfrm.get('flipV') == '1'
        off = xfrm.find(qn('a:off')); ext = xfrm.find(qn('a:ext'))
        x = int(off.get('x')); y = int(off.get('y'))
        w = int(ext.get('cx')); h = int(ext.get('cy'))
        bx, ex = (x + w, x) if flipH else (x, x + w)
        by, ey = (y + h, y) if flipV else (y, y + h)
        ln = el.find(qn('p:spPr')).find(qn('a:ln'))
        arrow = ln is not None and ln.find(qn('a:tailEnd')) is not None
        conns.append((emu_in(bx), emu_in(by), emu_in(ex), emu_in(ey), arrow))
    elif tag == 'sp':
        # text content + sizes
        txt = sp.text_frame.text if sp.has_text_frame else ''
        sizes = [int(s.get('sz')) for s in el.iter(qn('a:rPr')) if s.get('sz')]
        # detect named box
        node = txt.split('\n')[0].strip() if txt else ''
        sppr = el.find(qn('p:spPr'))
        xfrm = sppr.find(qn('a:xfrm')) if sppr is not None else None
        if node in ('A1','A2','A3','A4','A5','A6'):
            off = xfrm.find(qn('a:off')); ext = xfrm.find(qn('a:ext'))
            l = emu_in(int(off.get('x'))); t = emu_in(int(off.get('y')))
            boxes[node] = (l, t, l + emu_in(int(ext.get('cx'))), t + emu_in(int(ext.get('cy'))))
        # invisible ports: no text, tiny
        if xfrm is not None and not txt.strip():
            ext = xfrm.find(qn('a:ext')); off = xfrm.find(qn('a:off'))
            if ext is not None and emu_in(int(ext.get('cx'))) < 0.1:
                ports.append((emu_in(int(off.get('x'))) + emu_in(int(ext.get('cx')))/2,
                              emu_in(int(off.get('y'))) + emu_in(int(ext.get('cy')))/2))
        if txt.strip():
            texts.append((txt.replace('\n',' / '), sizes))

print("=== BOXES ===")
for n in sorted(boxes):
    l,t,r,b = boxes[n]; print(f"  {n}: x[{l:.2f},{r:.2f}] y[{t:.2f},{b:.2f}]")
print(f"\n#connectors={len(conns)}  #ports={len(ports)}  #textboxes={len(texts)}")

# ---- FONT AUDIT ----
bad_fonts = [(t,s) for t,s in texts for sz in s if sz < 1000]
print("\n=== FONT AUDIT (>=10pt) ===")
print("  FAIL:" if bad_fonts else "  PASS: all runs >= 10pt", bad_fonts[:5])

# ---- ORTHOGONALITY ----
diag = [c for c in conns if abs(c[0]-c[2])>0.005 and abs(c[1]-c[3])>0.005]
print("\n=== ORTHOGONALITY ===")
print("  PASS: all segments axis-aligned" if not diag else f"  FAIL diagonal: {diag}")

# ---- ENDPOINT-ON-FACE AUDIT ----
TOL = 0.06
def on_face(px, py):
    """Return (node, side) if (px,py) lies on a box perimeter face."""
    for n,(l,t,r,b) in boxes.items():
        if abs(px-l)<TOL and t-TOL<=py<=b+TOL: return (n,'L')
        if abs(px-r)<TOL and t-TOL<=py<=b+TOL: return (n,'R')
        if abs(py-t)<TOL and l-TOL<=px<=r+TOL: return (n,'T')
        if abs(py-b)<TOL and l-TOL<=px<=r+TOL: return (n,'B')
    return None

def inside(px,py):
    for n,(l,t,r,b) in boxes.items():
        if l+TOL<px<r-TOL and t+TOL<py<b-TOL: return n
    return None

print("\n=== ARROWHEAD ENDPOINT AUDIT (arrowed segments) ===")
arrow_ends = [c for c in conns if c[4]]
inside_hits = 0
face_hits = 0
boundary_hits = 0
for (bx,by,ex,ey,_) in arrow_ends:
    f = on_face(ex,ey)
    if inside(ex,ey): inside_hits += 1
    elif f: face_hits += 1
    elif ex>10.4 or ex<0.3: boundary_hits += 1
print(f"  arrowed segments: {len(arrow_ends)}")
print(f"  terminate on box face: {face_hits}")
print(f"  terminate at boundary (outputs): {boundary_hits}")
print(f"  arrowhead INSIDE a box (must be 0): {inside_hits}")

# ---- PORT-ON-FACE AUDIT ----
ports_on_face = sum(1 for (px,py) in ports if on_face(px,py))
print(f"\n=== PORTS ===\n  ports total {len(ports)}, on a box face: {ports_on_face}")

# perpendicular check: arrowed seg final approach must be perpendicular to face
print("\n=== PERPENDICULAR APPROACH (arrowed -> face) ===")
perp_ok=0; perp_bad=[]
for (bx,by,ex,ey,_) in arrow_ends:
    f=on_face(ex,ey)
    if not f:
        continue
    n,side=f
    horiz = abs(by-ey)<0.02   # last segment horizontal
    vert  = abs(bx-ex)<0.02
    if side in ('L','R') and horiz: perp_ok+=1
    elif side in ('T','B') and vert: perp_ok+=1
    else: perp_bad.append((n,side,round(bx,2),round(by,2),round(ex,2),round(ey,2)))
print(f"  perpendicular OK: {perp_ok}; problems: {perp_bad}")
