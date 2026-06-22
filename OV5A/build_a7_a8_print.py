#!/usr/bin/env python3
"""Build ATMOS_A7_A8_native_print_clean.pptx — letter-size (11x8.5 landscape),
government-print-ready IDEF0. Two slides:
  Slide 1: A7 External Reference (Execute Mission Adaptation) — external to ATMOS.
  Slide 2: A8 internal decomposition (Assess Outcomes and Learn).

Print standards enforced:
* slide 11.0 x 8.5 in, landscape
* no text below 10 pt
* function boxes contain ONLY node number + function name (no descriptions)
* native rectangles, editable text boxes, orthogonal elbow connector segments,
  native arrowheads, invisible connector-port shapes where >1 arrow lands on a
  box side
* arrowheads land exactly on outside perimeters; no connector through a box.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(11.0)
prs.slide_height = Inches(8.5)


def seg(sl, x1, y1, x2, y2, arrow=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.25)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(sl, points, arrow=True):
    """Orthogonal polyline (H/V segments only); arrowhead on final segment."""
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        assert abs(x1 - x2) < 1e-6 or abs(y1 - y2) < 1e-6, "non-orthogonal segment"
        seg(sl, x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def port(sl, x, y):
    """Tiny invisible connector-port shape centered exactly on a box perimeter."""
    sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x - 0.02), Inches(y - 0.02),
                             Inches(0.04), Inches(0.04))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def label(sl, x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    assert size >= 10, f"font below 10pt: {size} for {text!r}"
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0.5))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def fbox(sl, L, R, T, B, nid, name):
    """Function box: node number (line 1) + function name (line 2) only."""
    rect = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(L), Inches(T),
                               Inches(R - L), Inches(B - T))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(1.5); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = nid
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = name
    r2.font.size = Pt(10); r2.font.bold = False; r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    return rect


def frame(sl, L, T, R, B):
    fr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(L), Inches(T),
                             Inches(R - L), Inches(B - T))
    fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.0)
    fr.shadow.inherit = False


# ===================================================================
# SLIDE 1 — A7 EXTERNAL REFERENCE  (11 x 8.5)
# ===================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
FL, FR, FT, FB = 0.3, 10.7, 1.2, 8.1
frame(s1, FL, FT, FR, FB)
label(s1, 0.3, 0.04, 7.2, 0.3, "A7 — External Reference: Execute Mission Adaptation",
      size=12, bold=True)
label(s1, 7.5, 0.06, 3.2, 0.24, "Distribution Statement C", size=10, align=PP_ALIGN.RIGHT)
label(s1, 0.3, 0.40, 10.4, 0.22,
      "External C2 / mission-system activity — not an internal ATMOS A0 child decomposition",
      size=10, align=PP_ALIGN.CENTER)
# control labels (above frame top -> never crossed by drops that start at FT)
label(s1, 1.7, 0.66, 2.3, 0.46, "Command authority", size=10, align=PP_ALIGN.CENTER)
label(s1, 4.2, 0.66, 2.4, 0.46, "Rules of engagement", size=10, align=PP_ALIGN.CENTER)
label(s1, 6.8, 0.66, 2.4, 0.46, "Reporting policy", size=10, align=PP_ALIGN.CENTER)

# boxes (2x2)
A71 = (1.9, 3.8, 2.5, 3.45); A72 = (1.9, 3.8, 6.35, 7.3)
A73 = (5.6, 7.5, 2.5, 3.45); A74 = (5.6, 7.5, 6.35, 7.3)
fbox(s1, *A71, "A71", "Exchange Mission Adjustment Commands")
fbox(s1, *A72, "A72", "Publish Vertical Operations Execution Updates")
fbox(s1, *A73, "A73", "Publish UAS Tasking Updates")
fbox(s1, *A74, "A74", "Publish Coordination Status")

# node + external-boundary banner
label(s1, 7.6, 7.55, 3.0, 0.24, "Node: A7 External Reference", size=10, bold=True, align=PP_ALIGN.RIGHT)
label(s1, 7.7, 4.74, 2.95, 0.5, "EXTERNAL TO ATMOS SYSTEM BOUNDARY",
      size=11, bold=True, align=PP_ALIGN.CENTER)

# ---- inputs (left boundary) ----
seg(s1, FL, 2.95, 0.95, 2.95)
path(s1, [(0.95, 2.95), (1.9, 2.95)]); port(s1, 1.9, 2.95)                         # -> A71 left
path(s1, [(0.95, 2.95), (0.95, 4.0), (5.05, 4.0), (5.05, 2.9), (5.6, 2.9)]); port(s1, 5.6, 2.9)  # -> A73 left
path(s1, [(FL, 6.7), (1.9, 6.7)]); port(s1, 1.9, 6.7)                              # aircraft exec -> A72 left
path(s1, [(FL, 5.5), (4.9, 5.5), (4.9, 6.9), (5.6, 6.9)]); port(s1, 5.6, 6.9)      # coordination -> A74 left
label(s1, 0.33, 2.5, 1.5, 0.42, "External mission decisions", size=10)
label(s1, 0.33, 6.22, 1.5, 0.42, "Aircraft execution state", size=10)
label(s1, 0.33, 5.02, 1.5, 0.42, "Coordination state", size=10)

# ---- controls (top -> box top), ports for multi-arrow sides ----
seg(s1, 3.7, FT, 3.7, 1.7); seg(s1, 2.4, 1.7, 6.6, 1.7)                            # command authority bus
path(s1, [(2.4, 1.7), (2.4, 2.5)]); port(s1, 2.4, 2.5)                             # -> A71 top
path(s1, [(6.6, 1.7), (6.6, 2.5)]); port(s1, 6.6, 2.5)                             # -> A73 top
path(s1, [(3.0, FT), (3.0, 2.5)]); port(s1, 3.0, 2.5)                              # rules of engagement -> A71 top
seg(s1, 4.5, FT, 4.5, 5.7); seg(s1, 2.8, 5.7, 6.6, 5.7)                            # reporting policy bus
path(s1, [(2.8, 5.7), (2.8, 6.35)]); port(s1, 2.8, 6.35)                           # -> A72 top
path(s1, [(6.6, 5.7), (6.6, 6.35)]); port(s1, 6.6, 6.35)                           # -> A74 top

# ---- internal flows (labels clustered in the clear central pocket) ----
path(s1, [(3.8, 3.1), (5.6, 3.1)]); port(s1, 5.6, 3.1)                             # A7-F1 A71->A73 left
path(s1, [(6.7, 3.45), (6.7, 6.35)]); port(s1, 6.7, 6.35)                          # A7-F2 A73->A74 top
path(s1, [(3.8, 6.7), (5.6, 6.7)]); port(s1, 5.6, 6.7)                             # A7-F3 A72->A74 left
label(s1, 4.92, 4.42, 1.7, 0.34, "A7-F1 Mission decisions", size=10)
label(s1, 4.92, 4.80, 1.7, 0.34, "A7-F2 Tasking updates", size=10)
label(s1, 4.92, 5.18, 1.7, 0.34, "A7-F3 Execution updates", size=10)

# ---- outputs (right boundary) + labels in clear right-region pockets ----
path(s1, [(3.8, 2.7), (4.2, 2.7), (4.2, 4.3), (FR, 4.3)])                          # IER-19 from A71
path(s1, [(3.8, 6.95), (4.25, 6.95), (4.25, 5.95), (FR, 5.95)])                    # IER-20 from A72
path(s1, [(7.5, 2.85), (FR, 2.85)])                                               # IER-21 from A73
path(s1, [(7.5, 6.95), (FR, 6.95)])                                               # IER-22 from A74
label(s1, 7.75, 2.42, 2.9, 0.36, "UAS Tasking Updates (IER-21)", size=10)
label(s1, 7.75, 3.86, 2.9, 0.40, "Mission Adjustment Commands (IER-19)", size=10)
label(s1, 7.75, 5.50, 2.9, 0.40, "Vertical Ops Execution Updates (IER-20)", size=10)
label(s1, 7.75, 6.52, 2.9, 0.40, "Coordination Status Updates (IER-22)", size=10)

# ---- mechanisms (bottom -> box bottom), ports for multi-arrow sides ----
seg(s1, 4.8, FB, 4.8, 3.95); seg(s1, 2.4, 3.95, 6.6, 3.95)                         # C2 systems/TOC bus
path(s1, [(2.4, 3.95), (2.4, 3.45)]); port(s1, 2.4, 3.45)                          # -> A71 bottom
path(s1, [(6.6, 3.95), (6.6, 3.45)]); port(s1, 6.6, 3.45)                          # -> A73 bottom
path(s1, [(2.6, FB), (2.6, 7.3)]); port(s1, 2.6, 7.3)                              # aircraft avionics -> A72 bottom
seg(s1, 4.0, FB, 4.0, 7.65); seg(s1, 2.9, 7.65, 6.7, 7.65)                         # comms/DDS bus
path(s1, [(2.9, 7.65), (2.9, 7.3)]); port(s1, 2.9, 7.3)                            # -> A72 bottom
path(s1, [(6.7, 7.65), (6.7, 7.3)]); port(s1, 6.7, 7.3)                            # -> A74 bottom
label(s1, 1.55, 8.14, 1.7, 0.32, "Aircraft avionics", size=10, align=PP_ALIGN.CENTER)
label(s1, 3.35, 8.14, 2.7, 0.32, "Communications network / DDS", size=10, align=PP_ALIGN.CENTER)
label(s1, 6.2, 8.14, 2.0, 0.32, "C2 systems / TOC", size=10, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 2 — A8 INTERNAL DECOMPOSITION  (11 x 8.5)
# ===================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
FL2, FR2, FT2, FB2 = 0.3, 10.7, 1.05, 7.7
frame(s2, FL2, FT2, FR2, FB2)
label(s2, 0.3, 0.06, 7.5, 0.3, "A8 — Decompose Assess Outcomes and Learn", size=12, bold=True)
label(s2, 7.8, 0.08, 2.9, 0.24, "Distribution Statement C", size=10, align=PP_ALIGN.RIGHT)
label(s2, 9.2, 7.74, 1.45, 0.24, "Node: A8", size=10, bold=True, align=PP_ALIGN.RIGHT)

B81 = (2.2, 3.9, 3.4, 4.6); B82 = (4.75, 6.45, 3.4, 4.6); B83 = (7.3, 9.0, 3.4, 4.6)
cx = lambda b: (b[0] + b[1]) / 2
fbox(s2, *B81, "A81", "Compare Predicted vs Observed Conditions")
fbox(s2, *B82, "A82", "Update Models and Parameters")
fbox(s2, *B83, "A83", "Archive Lessons Learned")

# ---- inputs ----
path(s2, [(FL2, 3.75), (B81[0], 3.75)]); port(s2, B81[0], 3.75)                    # Fused COWP -> A81 left
path(s2, [(FL2, 4.30), (B81[0], 4.30)]); port(s2, B81[0], 4.30)                    # Observations -> A81 left
path(s2, [(FL2, 2.9), (4.5, 2.9), (4.5, 4.25), (B82[0], 4.25)]); port(s2, B82[0], 4.25)  # Archived Updates -> A82 left
label(s2, 0.33, 3.40, 1.8, 0.32, "Fused COWP State (IER-12)", size=10)
label(s2, 0.33, 4.37, 1.8, 0.32, "Observations (IER-01 / IER-02)", size=10)
label(s2, 0.33, 2.46, 2.0, 0.32, "Archived Updates (IER-24)", size=10)

# ---- controls ----
path(s2, [(cx(B81), FT2), (cx(B81), 3.4)]); port(s2, cx(B81), 3.4)
path(s2, [(cx(B82), FT2), (cx(B82), 3.4)]); port(s2, cx(B82), 3.4)
path(s2, [(cx(B83), FT2), (cx(B83), 3.4)]); port(s2, cx(B83), 3.4)
label(s2, 1.85, 0.62, 2.4, 0.3, "Assessment criteria", size=10, align=PP_ALIGN.CENTER)
label(s2, 4.4, 0.62, 2.4, 0.3, "Model governance", size=10, align=PP_ALIGN.CENTER)
label(s2, 6.95, 0.62, 2.5, 0.3, "Data retention / OPSEC", size=10, align=PP_ALIGN.CENTER)

# ---- internal flows (labels above boxes in clear band) ----
path(s2, [(B81[1], 3.9), (B82[0], 3.9)]); port(s2, B82[0], 3.9)                    # A8-F1 A81->A82 left
path(s2, [(B82[1], 4.0), (B83[0], 4.0)]); port(s2, B83[0], 4.0)                    # A8-F2 A82->A83 left
label(s2, 3.55, 2.16, 1.95, 0.32, "A8-F1 Assessment Outputs", size=10, align=PP_ALIGN.CENTER)
label(s2, 5.95, 2.96, 1.95, 0.3, "A8-F2 Updated Models", size=10, align=PP_ALIGN.CENTER)

# ---- outputs ----
path(s2, [(B81[1], 3.62), (4.45, 3.62), (4.45, 2.55), (FR2, 2.55)])                # IER-23 from A81 -> boundary
path(s2, [(B83[1], 4.2), (FR2, 4.2)])                                             # IER-24 from A83 -> boundary
label(s2, 8.5, 2.05, 2.15, 0.46, "Predicted vs Observed Assessment (IER-23)", size=10)
label(s2, 9.05, 3.34, 1.6, 0.8, "Archived Lessons Learned & Model Updates (IER-24)", size=10, bold=True)

# ---- mechanisms ----
path(s2, [(cx(B81), FB2), (cx(B81), 4.6)]); port(s2, cx(B81), 4.6)
path(s2, [(cx(B82), FB2), (cx(B82), 4.6)]); port(s2, cx(B82), 4.6)
path(s2, [(cx(B83), FB2), (cx(B83), 4.6)]); port(s2, cx(B83), 4.6)
label(s2, 1.85, 7.74, 2.4, 0.3, "ATMOS analytics", size=10, align=PP_ALIGN.CENTER)
label(s2, 4.4, 7.74, 2.4, 0.3, "Model management tools", size=10, align=PP_ALIGN.CENTER)
label(s2, 6.95, 7.74, 2.5, 0.3, "TruWeather Data Warehouse", size=10, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A7_A8_native_print_clean.pptx")
print("saved ATMOS_A7_A8_native_print_clean.pptx")
