#!/usr/bin/env python3
"""
Generate ATMOS_A8_IDEF0_native_vNext.pptx
A8 LIFECYCLE/SUPPORT IDEF0: chain A81->A82->A83. 11x8.5 landscape, reference
frame style, bare 'A8' node id, scope note. >=10pt.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK=RGBColor(0,0,0); WHITE=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(11.0); prs.slide_height=Inches(8.5)
slide=prs.slides.add_slide(prs.slide_layouts[6]); shapes=slide.shapes
def IN(v): return Inches(v)

def add_box(node,name,l,t,w,h):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=WHITE
    sp.line.color.rgb=BLACK; sp.line.width=Pt(1.5); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0=tf.paragraphs[0]; p0.alignment=PP_ALIGN.CENTER
    r0=p0.add_run(); r0.text=node; r0.font.size=Pt(12); r0.font.bold=True; r0.font.color.rgb=BLACK
    p1=tf.add_paragraph(); p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text=name; r1.font.size=Pt(10); r1.font.color.rgb=BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb=shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb=WHITE; tb.line.fill.background()
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=BLACK
    return tb

PORT=0.035
def add_port(x,y):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit=False

def _arrow(c):
    ln=c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c=shapes.add_connector(MSO_CONNECTOR.ELBOW,IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb=BLACK; c.line.width=Pt(1.25); c.shadow.inherit=False
    if arrow:_arrow(c)
    return c

def path(pts,src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code,x,y):
    add_label(code,x-0.33,y,0.66,0.24,10,bold=True,align=PP_ALIGN.CENTER,fill_white=True)

# ---------------- frame ----------------
bd=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb=BLACK; bd.line.width=Pt(1.25); bd.shadow.inherit=False
add_label('A8 — Decompose Assess Outcomes and Learn',
          0.35,0.60,10.30,0.32,15,bold=True,align=PP_ALIGN.CENTER)
add_label('LIFECYCLE / SUPPORT ACTIVITY',0.55,0.96,3.5,0.28,11,bold=True)
nb=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(9.35),IN(7.62),IN(1.20),IN(0.30))
nb.fill.background(); nb.line.color.rgb=BLACK; nb.line.width=Pt(1.0); nb.shadow.inherit=False
nt=nb.text_frame; nt.vertical_anchor=MSO_ANCHOR.MIDDLE; nt.paragraphs[0].alignment=PP_ALIGN.CENTER
nr=nt.paragraphs[0].add_run(); nr.text='A8'; nr.font.size=Pt(12); nr.font.bold=True; nr.font.color.rgb=BLACK

# ---------------- boxes ----------------
add_box('A81','Compare Predicted vs Observed Conditions',1.30,3.50,2.35,0.90)
add_box('A82','Update Models and Parameters',4.30,3.50,2.25,0.90)
add_box('A83','Archive Lessons Learned',7.30,3.50,2.25,0.90)
CY=3.95; TOPF=3.50; BOTF=4.40; RBX=10.55

# ---------------- INPUTS (left) ----------------
path([(0.50,3.72),(1.30,3.72)])                               # Fused COWP State -> A81 left
path([(0.50,4.12),(1.30,4.12)])                               # Observations -> A81 left
add_label('Fused COWP State (IER-12)',0.40,3.08,1.85,0.36,10,fill_white=True)
add_label('Observations (IER-01 / IER-02)',0.40,4.46,1.85,0.36,10,fill_white=True)
# Prior Archived Updates / Model Baselines -> A82 left (around A81)
path([(0.50,2.95),(4.10,2.95),(4.10,3.72),(4.30,3.72)])
add_label('Prior Archived Updates / Model Baselines (IER-24 reference)',
          1.60,2.56,2.30,0.34,10,fill_white=True)

# ---------------- CONTROLS (top) ----------------
path([(2.475,1.55),(2.475,TOPF)])                             # C81 -> A81
path([(5.425,1.55),(5.425,TOPF)])                             # C82 -> A82
path([(8.425,1.55),(8.425,TOPF)])                             # C83 -> A83
add_label('C81 Assessment criteria; time-alignment rules',1.25,1.00,2.30,0.42,10)
add_label('C82 Model governance; validation / approval process',4.15,1.00,2.45,0.42,10)
add_label('C83 Data retention; OPSEC / classification; sharing constraints',7.05,1.00,2.55,0.55,10)

# ---------------- INTERNAL FLOWS ----------------
path([(3.65,4.12),(4.30,4.12)],src_port=True)                 # A8-F1 A81->A82 left
path([(6.55,CY),(7.30,CY)],src_port=True)                     # A8-F2 A82->A83 left
add_label('A8-F1 Assessment Outputs',2.95,4.85,1.85,0.34,10,fill_white=True)
add_label('A8-F2 Updated Model Parameters and Configuration Baselines',
          5.85,4.85,2.05,0.50,10,fill_white=True)

# ---------------- OUTPUTS (right) ----------------
# IER-23 from A81 right, routed above the chain to the boundary
path([(3.65,3.72),(3.95,3.72),(3.95,3.12),(RBX,3.12)],src_port=True)
# IER-24 from A83 right to boundary
path([(9.55,CY),(RBX,CY)],src_port=True)
add_label('Predicted vs Observed Weather Assessment (IER-23)',8.45,2.66,2.05,0.42,10,fill_white=True)
add_label('Archived Lessons Learned & Model Updates (IER-24)',8.45,3.30,2.05,0.42,10,fill_white=True)

# ---------------- MECHANISM-LANE MATRIX ----------------
BUS=6.20
path([(2.475,BUS),(2.475,BOTF)]); mtag('M81',2.475,4.50)
path([(5.425,BUS),(5.425,BOTF)]); mtag('M82',5.425,4.50)
path([(8.425,BUS),(8.425,BOTF)]); mtag('M83',8.425,4.50)
add_label('M81 ATMOS analytics; TOC analysis tooling; storage',0.55,6.55,2.55,0.55,10,fill_white=True)
add_label('M82 Model management tools; ABLE-LBM configuration management',3.85,6.50,2.75,0.70,10,fill_white=True)
add_label('M83 TruWeather Data Warehouse; ingestion pipeline',7.00,6.50,2.30,0.70,10,fill_white=True)

prs.save('ATMOS_A8_IDEF0_native_vNext.pptx')
print('saved; shapes:',len(slide.shapes._spTree))
