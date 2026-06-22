#!/usr/bin/env python3
"""Build ATMOS_A2_native_clean.pptx — native, editable IDEF0 A2 decomposition.

Authoritative model (inline from task): parent A2 'Generate Local Weather State'
-> A21, A22, A23. Internal flows A2-F1 (A21->A22), A2-F2 (A22->A23) only.
A24/TBD NOT drawn. COWP NOT drawn. Parent output F2 (IER-05) exits from A23.
Coordinate-placed right-angle connector paths; the terminal segment carries the
arrowhead exactly on each box's outside-perimeter anchor. No glue, no images,
no whole-diagram group, no connector grouped with a box.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    """Draw a right-angle polyline; arrowhead only on the final segment."""
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, desc, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(9); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.32, 0.45, 0.26, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect


# -------- frame: x 0.5..21.5, y 1.1..11.8 --------
FL, FT, FR, FB = 0.5, 1.1, 21.5, 11.8
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT),
                            Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# -------- title / marking / node --------
label(6.0, 0.12, 12.0, 0.42, "A2 — Decompose Generate Local Weather State",
      size=18, bold=True, align=PP_ALIGN.CENTER)
label(17.9, 0.16, 3.5, 0.3, "Distribution Statement C", size=12, align=PP_ALIGN.RIGHT)
label(19.5, 11.45, 1.9, 0.3, "Node: A2", size=12, bold=True, align=PP_ALIGN.RIGHT)

# -------- activity boxes: horizontal chain A21 -> A22 -> A23 --------
# A21: x 2.5..5.9  | A22: x 9.0..12.4 | A23: x 15.5..18.9 ; all y 4.2..5.8 (mid y=5.0)
BY1, BY2 = 4.2, 5.8
A21 = (2.5, 5.9); A22 = (9.0, 12.4); A23 = (15.5, 18.9)
box(A21[0], BY1, A21[1] - A21[0], BY2 - BY1,
    "Prepare Model Inputs and Boundary Conditions",
    "Format QC'd observations and cached context into model-ready inputs.", "A21")
box(A22[0], BY1, A22[1] - A22[0], BY2 - BY1,
    "Execute Local Micro-Weather Estimation",
    "Run ABLE-LBM or reduced models to generate local micro-weather state estimates.", "A22")
box(A23[0], BY1, A23[1] - A23[0], BY2 - BY1,
    "Derive Aviation-Relevant Weather Fields",
    "Transform model outputs into aviation-relevant weather fields and representations.", "A23")
MY = 5.0  # box mid-height (input/flow rail)

# -------- parent input F1 (LEFT boundary -> A21 left perimeter) --------
path([(FL, MY), (A21[0], MY)])
label(0.55, 4.28, 1.9, 0.62,
      "F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)", size=9)

# -------- internal flows A2-F1, A2-F2 --------
path([(A21[1], MY), (A22[0], MY)])   # A2-F1: A21 right -> A22 left
path([(A22[1], MY), (A23[0], MY)])   # A2-F2: A22 right -> A23 left
label(5.95, 4.34, 3.0, 0.6,
      "A2-F1 Model Inputs & Boundary Conditions (IER-04)", size=9, align=PP_ALIGN.CENTER)
label(12.45, 4.34, 3.0, 0.6,
      "A2-F2 Local Micro-Weather State Estimate / Model Output", size=9, align=PP_ALIGN.CENTER)

# -------- parent output F2 (A23 right perimeter -> RIGHT boundary) --------
path([(A23[1], MY), (FR, MY)])
label(19.0, 4.28, 2.4, 0.6,
      "F2 Local Micro-Weather State Estimate (IER-05)", size=9)
label(21.6, 4.82, 2.3, 0.5, "To A3 Quantify Uncertainty", size=9, bold=True)

# -------- controls (TOP boundary -> box top perimeter) --------
C21x, C22x, C23x = 4.2, 10.7, 17.2   # box center x
path([(C21x, FT), (C21x, BY1)])
path([(C22x, FT), (C22x, BY1)])
path([(C23x, FT), (C23x, BY1)])
label(C21x - 1.6, 0.58, 3.2, 0.48,
      "C21 Model configuration; resource limits; update cadence",
      size=9, align=PP_ALIGN.CENTER)
label(C22x - 1.6, 0.58, 3.2, 0.48,
      "C22 Execution timing; numerical stability constraints; compute/power constraints",
      size=9, align=PP_ALIGN.CENTER)
label(C23x - 1.6, 0.58, 3.2, 0.48,
      "C23 Field derivation mappings; aviation relevance rules; reporting format",
      size=9, align=PP_ALIGN.CENTER)

# -------- mechanisms (BOTTOM boundary -> box bottom perimeter) --------
# Distinct bottom-edge entry points per box so risers never overlap.
# A21 bottom (2.5..5.9): M1=3.0  M21=4.0  M22=5.0
# A22 bottom (9.0..12.4): M1=9.5 M2=10.4  M22=11.4
# A23 bottom (15.5..18.9): M1=16.0 M21=17.0 M22=18.0
# Stacked distribution buses below the boxes; one trunk each from bottom boundary.

# M1 -> A21, A22, A23  (bus y=7.0, trunk x=4.5)
M1y = 7.0
seg(4.5, FB, 4.5, M1y)                 # trunk from bottom boundary
seg(3.0, M1y, 16.0, M1y)               # distribution bus
path([(3.0, M1y), (3.0, BY2)])         # riser -> A21 bottom
path([(9.5, M1y), (9.5, BY2)])         # riser -> A22 bottom
path([(16.0, M1y), (16.0, BY2)])       # riser -> A23 bottom

# M2 -> A22 only  (direct trunk+riser, x=10.4)
path([(10.4, FB), (10.4, BY2)])        # bottom boundary -> A22 bottom

# M21 -> A21, A23  (bus y=8.0, trunk x=8.0)
M21y = 8.0
seg(8.0, FB, 8.0, M21y)                # trunk from bottom boundary
seg(4.0, M21y, 17.0, M21y)             # distribution bus
path([(4.0, M21y), (4.0, BY2)])        # riser -> A21 bottom
path([(17.0, M21y), (17.0, BY2)])      # riser -> A23 bottom

# M22 -> A21, A22, A23  (bus y=9.0, trunk x=13.0)
M22y = 9.0
seg(13.0, FB, 13.0, M22y)              # trunk from bottom boundary
seg(5.0, M22y, 18.0, M22y)             # distribution bus
path([(5.0, M22y), (5.0, BY2)])        # riser -> A21 bottom
path([(11.4, M22y), (11.4, BY2)])      # riser -> A22 bottom
path([(18.0, M22y), (18.0, BY2)])      # riser -> A23 bottom

# mechanism labels (below bottom boundary, near each trunk)
label(3.25, 11.92, 2.6, 0.5, "M1 Platforms hosting ATMOS node compute",
      size=9, align=PP_ALIGN.CENTER)
label(6.75, 11.92, 2.6, 0.6, "M21 ATMOS preprocessing and post-processing software",
      size=9, align=PP_ALIGN.CENTER)
label(9.15, 12.5, 2.5, 0.4, "M2 ABLE-LBM / reduced models",
      size=9, align=PP_ALIGN.CENTER)
label(11.8, 11.92, 2.5, 0.5, "M22 Local compute/storage",
      size=9, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A2_native_clean.pptx")
print("saved ATMOS_A2_native_clean.pptx; shapes:", len(slide.shapes))
