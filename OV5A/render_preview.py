#!/usr/bin/env python3
"""Render ATMOS_A0_IDEF0_native_vNext.pptx to PNG for visual self-audit.
Reads native shape geometry and draws boxes, connectors, arrowheads, labels."""
import sys
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrow
from pptx import Presentation
from pptx.oxml.ns import qn

EMU = 914400.0
fn = sys.argv[1] if len(sys.argv) > 1 else 'ATMOS_A0_IDEF0_native_vNext.pptx'
prs = Presentation(fn)
W = prs.slide_width / EMU; H = prs.slide_height / EMU
slide = prs.slides[0]

fig, ax = plt.subplots(figsize=(W, H), dpi=150)
ax.set_xlim(0, W); ax.set_ylim(0, H)
ax.invert_yaxis()                # PowerPoint y grows downward
ax.axis('off')

def emu(v): return v / EMU

for sp in slide.shapes:
    el = sp._element
    tag = el.tag.split('}')[-1]
    sppr = el.find(qn('p:spPr'))
    xfrm = sppr.find(qn('a:xfrm')) if sppr is not None else None
    if tag == 'cxnSp':
        flipH = xfrm.get('flipH') == '1'; flipV = xfrm.get('flipV') == '1'
        off = xfrm.find(qn('a:off')); ext = xfrm.find(qn('a:ext'))
        x = emu(int(off.get('x'))); y = emu(int(off.get('y')))
        w = emu(int(ext.get('cx'))); h = emu(int(ext.get('cy')))
        bx, ex = (x+w, x) if flipH else (x, x+w)
        by, ey = (y+h, y) if flipV else (y, y+h)
        ln = sppr.find(qn('a:ln'))
        arrow = ln is not None and ln.find(qn('a:tailEnd')) is not None
        ax.plot([bx, ex], [by, ey], color='black', lw=0.9, solid_capstyle='round')
        if arrow:
            dx, dy = ex-bx, ey-by
            n = (dx*dx+dy*dy)**0.5 or 1
            ux, uy = dx/n, dy/n
            ax.annotate('', xy=(ex, ey), xytext=(ex-ux*0.12, ey-uy*0.12),
                        arrowprops=dict(arrowstyle='-|>', color='black', lw=0.9))
        continue
    if xfrm is None:
        continue
    off = xfrm.find(qn('a:off')); ext = xfrm.find(qn('a:ext'))
    x = emu(int(off.get('x'))); y = emu(int(off.get('y')))
    w = emu(int(ext.get('cx'))); h = emu(int(ext.get('cy')))
    txt = sp.text_frame.text if sp.has_text_frame else ''
    # fill / line detection
    has_line = sppr.find(qn('a:ln')) is not None and \
        sppr.find(qn('a:ln')).find(qn('a:noFill')) is None
    if tag == 'sp':
        node = txt.split('\n')[0].strip()
        try:
            is_box = str(sp.fill.type) == 'SOLID (1)'
        except Exception:
            is_box = False
        is_box = (is_box or node in ('A1','A2','A3','A4','A5','A6')) and has_line
        is_frame = (w > 9 and h > 6)
        if has_line:    # only true boxes / frame / node-cell carry a visible border
            ax.add_patch(Rectangle((x, y), w, h, fill=False,
                                   edgecolor='black', lw=1.2 if is_box or is_frame else 0.8))
        if txt.strip():
            import textwrap
            fs = 7 if is_box else 5.5
            cpl = max(4, int(w * 72 / (fs * 0.62)))   # est chars per line
            wrapped = []
            for ln_ in txt.split('\n'):
                wrapped += textwrap.wrap(ln_, cpl) or ['']
            disp = '\n'.join(wrapped)
            centered = is_box or (w > 9)              # boxes + title centered
            ax.text(x + w/2 if centered else x + 0.02,
                    y + h/2,
                    disp, fontsize=fs, ha='center' if centered else 'left',
                    va='center',
                    fontweight='bold' if is_box else 'normal')

ax.add_patch(Rectangle((0, 0), W, H, fill=False, edgecolor='0.8', lw=0.5))
plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
out = fn.replace('.pptx', '_preview.png')
plt.savefig(out, dpi=150)
print('wrote', out)
