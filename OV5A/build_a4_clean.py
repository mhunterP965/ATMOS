#!/usr/bin/env python3
"""Build ATMOS_A4_native_clean.pptx — native, editable IDEF0 A4 decomposition.

Authoritative model (inline from task): parent A4 'Federate and Maintain
Federated Weather Context' -> A41, A42, A43, A44 (federated, NOT serial).

Layout (federation shape; A43 central, A44 below A43 with feedback loop):
    A41 (upper-left)              A43 (center-right) --F4--> right boundary
                       >--converge-->   |  ^
    A42 (lower-left)                     v  | (A4-F3 down / A4-F4 cached up)
                                       A44 (below A43)

Stacked columns force "around" routing: controls to the lower box of a column
route through the clear inter-row band; mechanisms to the upper box likewise.
The middle gap (x 6.0..11.0) and the inter-row bands are box-free highways.

A41 does NOT feed A42. F3 input branches to A41 and A43. Peer publications feed
A42. A41->A43 (A4-F1), A42->A43 (A4-F2), A43->A44 (A4-F3), A44->A43 (A4-F4).
Parent output F4 originates from A43 only. COWP appears only in the F4 label.

Coordinate-placed right-angle connector paths; arrowhead on the terminal segment,
exactly on each box's outside-perimeter anchor. No glue, no images, no group.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    """Right-angle polyline; arrowhead only on the final segment."""
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, desc, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(8.5); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.30, 0.45, 0.24, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect


# ---------------- frame ----------------
FL, FT, FR, FB = 0.5, 1.1, 21.5, 11.8
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(FL), Inches(FT),
                            Inches(FR - FL), Inches(FB - FT))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# ---------------- title / marking / node ----------------
label(4.5, 0.10, 15.0, 0.30, "A4 — Decompose Federate and Maintain Federated Weather Context",
      size=16, bold=True, align=PP_ALIGN.CENTER)
label(18.0, 0.10, 3.4, 0.26, "Distribution Statement C", size=11, align=PP_ALIGN.RIGHT)
label(19.5, 11.45, 1.9, 0.3, "Node: A4", size=12, bold=True, align=PP_ALIGN.RIGHT)

# ---------------- activity boxes ----------------
# A41 upper-left, A42 lower-left (same x column); A43 center-right; A44 below A43.
A41 = dict(L=2.5, R=6.0, T=3.0, B=4.5)
A42 = dict(L=2.5, R=6.0, T=7.0, B=8.5)
A43 = dict(L=11.0, R=15.0, T=4.6, B=6.6)
A44 = dict(L=11.0, R=15.0, T=8.4, B=9.9)
box(A41['L'], A41['T'], A41['R']-A41['L'], A41['B']-A41['T'],
    "Publish Local Weather State",
    "Publish locally generated weather state and uncertainty to the network via DDS.", "A41")
box(A42['L'], A42['T'], A42['R']-A42['L'], A42['B']-A42['T'],
    "Subscribe to Peer Weather States",
    "Receive weather states from peer nodes and optional reachback when connectivity permits.", "A42")
box(A43['L'], A43['T'], A43['R']-A43['L'], A43['B']-A43['T'],
    "Stitch and Blend Weather States",
    "Reconcile local and peer states into a fused, internally consistent shared context instance.", "A43")
box(A44['L'], A44['T'], A44['R']-A44['L'], A44['B']-A44['T'],
    "Maintain Degraded or Disconnected Operation",
    "Maintain cached weather states and continue operation under degraded or disconnected communications.", "A44")

# ======================= INPUTS =======================
# F3 -> branch to A41 (left) and A43 (left). Enters left boundary.
seg(FL, 3.6, 1.6, 3.6)                              # F3 trunk from boundary
path([(1.6, 3.6), (2.5, 3.6)])                      # F3 -> A41 left perimeter
path([(1.6, 3.6), (1.6, 5.4), (A43['L'], 5.4)])     # F3 -> A43 left perimeter
label(0.55, 2.86, 2.4, 0.6, "F3 Confidence Bounds & Risk Envelopes (IER-09)", size=8)

# Peer Weather State Publications -> A42 left
path([(FL, 7.6), (2.5, 7.6)])
label(0.55, 6.86, 2.4, 0.6, "Peer Weather State Publications (IER-10 / IER-11)", size=8)

# ======================= INTERNAL FLOWS =======================
# A4-F1: A41 right -> A43 left  (label placed in lane-free zone above the gap)
path([(6.0, 3.8), (8.0, 3.8), (8.0, 5.0), (A43['L'], 5.0)])
label(7.05, 2.50, 2.5, 0.54, "A4-F1 Local Weather State Publication / Local State Available (IER-10)", size=8)
# A4-F2: A42 right -> A43 left  (label placed in lane-free zone below the gap)
path([(6.0, 7.6), (8.6, 7.6), (8.6, 5.8), (A43['L'], 5.8)])
label(7.05, 8.58, 2.5, 0.54, "A4-F2 Peer Weather State Subscription (IER-11)", size=8)
# A4-F3: A43 bottom -> A44 top  (fused state down)
path([(12.5, A43['B']), (12.5, A44['T'])])
# A4-F4: A44 top -> A43 bottom  (cached state up, feedback)
path([(13.5, A44['T']), (13.5, A43['B'])])
label(15.15, 6.85, 2.5, 0.42, "A4-F3 Fused Weather Context State (IER-12)", size=8)
label(15.15, 7.55, 2.5, 0.42, "A4-F4 Cached Weather State", size=8)

# ======================= OUTPUT =======================
# F4: A43 right -> right boundary
path([(A43['R'], 5.6), (FR, 5.6)])
label(15.25, 4.92, 3.4, 0.42, "F4 Federated Weather Context / COWP State", size=8)
label(21.6, 5.18, 2.3, 0.9, "To A5 Detect Threshold Crossings and A6 Produce Mission-Tailored Weather Context",
      size=8, bold=True)

# ======================= CONTROLS (top boundary) =======================
# C3 -> A41, A42, A44  (top bus at y=1.4)
seg(5.0, FT, 5.0, 1.4)                               # C3 trunk from boundary
seg(3.5, 1.4, 10.2, 1.4)                             # C3 bus
path([(3.5, 1.4), (3.5, A41['T'])])                  # C3 -> A41 top
path([(6.6, 1.4), (6.6, 5.6), (3.0, 5.6), (3.0, A42['T'])])    # C3 -> A42 top (around A41)
path([(10.2, 1.4), (10.2, 7.4), (11.6, 7.4), (11.6, A44['T'])])# C3 -> A44 top (around A43)
# C2 -> A42, A43, A44  (top bus at y=1.7)
seg(9.0, FT, 9.0, 1.7)                               # C2 trunk from boundary
seg(6.2, 1.7, 11.8, 1.7)                             # C2 bus
path([(11.8, 1.7), (11.8, A43['T'])])                # C2 -> A43 top
path([(6.2, 1.7), (6.2, 6.0), (3.6, 6.0), (3.6, A42['T'])])    # C2 -> A42 top (around A41)
path([(10.6, 1.7), (10.6, 7.0), (12.0, 7.0), (12.0, A44['T'])])# C2 -> A44 top (around A43)
# C41 -> A43 (single)
path([(14.2, FT), (14.2, A43['T'])])
# C42 -> A44 (single, around A43)
path([(9.8, FT), (9.8, 7.8), (14.2, 7.8), (14.2, A44['T'])])
# control labels (two rows above frame top)
label(2.2, 0.40, 3.6, 0.30, "C3 Network availability and DDS QoS policies", size=8, align=PP_ALIGN.CENTER)
label(6.3, 0.40, 3.2, 0.30, "C2 OPSEC / classification guidance", size=8, align=PP_ALIGN.CENTER)
label(11.6, 0.40, 4.2, 0.30, "C41 Fusion heuristics; data age limits; prioritization rules", size=8, align=PP_ALIGN.CENTER)
label(6.2, 0.76, 4.4, 0.30, "C42 Cache expiration; reconnection behavior; storage limits", size=8, align=PP_ALIGN.CENTER)

# ======================= MECHANISMS (bottom boundary) =======================
# M1 -> A41, A42, A43, A44  (bottom bus at y=11.3)
seg(8.0, FB, 8.0, 11.3)                              # M1 trunk from boundary
seg(5.0, 11.3, 12.0, 11.3)                           # M1 bus
path([(5.0, 11.3), (5.0, A42['B'])])                 # M1 -> A42 bottom
path([(6.4, 11.3), (6.4, 5.3), (4.6, 5.3), (4.6, A41['B'])])   # M1 -> A41 bottom (around A42)
path([(10.4, 11.3), (10.4, 7.2), (11.5, 7.2), (11.5, A43['B'])])# M1 -> A43 bottom (around A44)
path([(12.0, 11.3), (12.0, A44['B'])])               # M1 -> A44 bottom
# M3 -> A41, A42  (bottom bus at y=10.7)
seg(6.2, FB, 6.2, 10.7)                              # M3 trunk from boundary
seg(5.5, 10.7, 7.0, 10.7)                            # M3 bus
path([(5.5, 10.7), (5.5, A42['B'])])                 # M3 -> A42 bottom
path([(7.0, 10.7), (7.0, 5.0), (4.0, 5.0), (4.0, A41['B'])])   # M3 -> A41 bottom (around A42)
# M41 -> A43 (single, around A44)
path([(9.6, FB), (9.6, 6.9), (14.2, 6.9), (14.2, A43['B'])])
# M42 -> A44 (single)
path([(13.0, FB), (13.0, A44['B'])])
# mechanism labels (below frame bottom)
label(5.6, 11.95, 3.0, 0.30, "M1 Platforms hosting ATMOS node compute", size=8, align=PP_ALIGN.CENTER)
label(8.9, 11.95, 2.6, 0.30, "M41 ATMOS federation logic", size=8, align=PP_ALIGN.CENTER)
label(11.7, 11.95, 2.8, 0.30, "M42 Local storage / cache manager", size=8, align=PP_ALIGN.CENTER)
label(4.4, 12.35, 3.4, 0.30, "M3 DDS middleware and communications links", size=8, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A4_native_clean.pptx")
print("saved ATMOS_A4_native_clean.pptx; shapes:", len(slide.shapes))
