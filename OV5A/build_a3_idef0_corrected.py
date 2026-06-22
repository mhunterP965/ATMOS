#!/usr/bin/env python3
"""
Generate ATMOS_A3_IDEF0_native_vNext.pptx

A3 IDEF0 decomposition (A31 -> A32 -> A33). 11.0 x 8.5 landscape, reference
frame style. F2 input branches to A32 and A33; A3-F3 combined flow from
A31+A32 into A33. Auditable mechanism-lane matrix. >= 10 pt fonts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0,0,0); WHITE = RGBColor(0xFF,0xFF,0xFF)
prs = Presentation(); prs.slide_width = Inches(11.0); prs.slide_height = Inches(8.5)
slide = prs.slides.add_slide(prs.slide_layouts[6]); shapes = slide.shapes
def IN(v): return Inches(v)

def add_box(node, name, l, t, w, h):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = BLACK; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in ('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = node; r0.font.size = Pt(12); r0.font.bold = True
    r0.font.color.rgb = BLACK
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = name; r1.font.size = Pt(10); r1.font.color.rgb = BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb = shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb = WHITE; tb.line.fill.background()
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK
    return tb

PORT = 0.035
def add_port(x,y):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit = False

def _arrow(c):
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c = shapes.add_connector(MSO_CONNECTOR.ELBOW, IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.25); c.shadow.inherit = False
    if arrow: _arrow(c)
    return c

def path(pts, src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code, x, y):
    add_label(code, x-0.31, y, 0.62, 0.24, 10, bold=True, align=PP_ALIGN.CENTER,
              fill_white=True)

# ----------------------------------------------------------------- frame
bd = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb = BLACK; bd.line.width = Pt(1.25); bd.shadow.inherit = False
add_label('A3 — Decompose Quantify Uncertainty',
          0.35,0.60,10.30,0.34,15,bold=True,align=PP_ALIGN.CENTER)
nb = shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(9.10),IN(7.62),IN(1.45),IN(0.30))
nb.fill.background(); nb.line.color.rgb = BLACK; nb.line.width = Pt(1.0); nb.shadow.inherit = False
nt = nb.text_frame; nt.vertical_anchor = MSO_ANCHOR.MIDDLE
nt.paragraphs[0].alignment = PP_ALIGN.CENTER
nr = nt.paragraphs[0].add_run(); nr.text = 'Node: A3'; nr.font.size = Pt(11)
nr.font.bold = True; nr.font.color.rgb = BLACK

# ----------------------------------------------------------------- boxes
add_box('A31','Estimate Observation Uncertainty',1.25,3.55,2.25,0.90)
add_box('A32','Estimate Model and State Uncertainty',4.15,3.55,2.35,0.90)
add_box('A33','Package Confidence Bounds and Risk Envelopes',7.25,3.55,2.35,0.90)

CY = 4.00; TOPF = 3.55; BOTF = 4.45; RBX = 10.45

# ============================ CONTROLS (top) ============================
path([(2.375,1.55),(2.375,TOPF)])                  # C31 -> A31 top
path([(5.325,1.55),(5.325,TOPF)])                  # C32 -> A32 top
path([(8.375,1.55),(8.375,TOPF)])                  # C33 -> A33 top
add_label('C31 Sensor uncertainty models; calibration/state-of-health policies',
          1.20,0.98,2.30,0.55,10)
add_label('C32 Error propagation rules; model tuning parameters',
          3.95,1.00,2.45,0.50,10)
add_label('C33 Confidence representation standards; formatting constraints',
          6.95,1.00,2.55,0.55,10)

# ============================ INPUT F2 (inherited A2 output -> A3 left inputs) ============================
# Clean left-side input bus from the boundary with separate horizontal stubs into
# A32 and A33; stub x-positions kept clear of the A3-F3 lane so F2 stays distinct.
seg(0.50,2.85,7.00,2.85)                                          # F2 input bus (from left boundary)
path([(3.58,2.85),(3.58,3.80),(4.15,3.80)])                       # F2 -> A32 left (dedicated port)
path([(7.00,2.85),(7.00,3.70),(7.25,3.70)])                       # F2 -> A33 left (dedicated port)
add_label('F2 Local Micro-Weather State Estimate (IER-05)',
          0.42,2.14,1.95,0.34,10, bold=True, fill_white=True)
add_label('From A2 Generate Local Weather State',
          0.45,2.49,1.95,0.30,10, fill_white=True)

# ============================ INTERNAL FLOWS ============================
path([(3.50,CY),(4.15,CY)], src_port=True)                        # A3-F1 A31->A32
path([(6.50,CY),(7.25,CY)], src_port=True)                        # A3-F2 A32->A33
# A3-F3 combined: A31 + A32 right outputs -> A33 left (single arrow at A33)
path([(3.50,3.70),(3.65,3.70),(3.65,3.15),(6.95,3.15),(6.95,4.25),(7.25,4.25)],
     src_port=True)                                               # A31 contribution + lane + A33
seg(6.50,3.70,6.70,3.70); seg(6.70,3.70,6.70,3.15); add_port(6.50,3.70)  # A32 joins lane
add_label('A3-F1 Observation Uncertainty Metrics (IER-07)',
          2.85,4.60,1.95,0.42,10, fill_white=True)
add_label('A3-F2 Model / State Uncertainty Estimates (IER-08)',
          6.52,4.60,1.95,0.42,10, fill_white=True)
add_label('A3-F3 Combined Uncertainty Inputs',
          4.05,2.78,1.95,0.34,10, fill_white=True)

# ============================ OUTPUT F3 (right) ============================
path([(9.60,CY),(RBX,CY)], src_port=True)                         # F3 -> right boundary
add_label('F3 Confidence Bounds & Risk Envelopes (IER-09)',
          9.05,3.00,1.58,0.50,10, fill_white=True)
add_label('To A4 Federate and Maintain Federated Weather Context',
          9.62,4.06,1.02,0.70,10, fill_white=True)

# ============================ MECHANISM-LANE MATRIX ============================
# M1 spine (all three) at lower lane; single-target mechanisms as short stubs on
# an upper lane so nothing crosses (different x columns).  Per-branch code tags.
SING = 6.10; M1y = 6.50
# bottom ports
A31_M1, A31_M31 = 1.95, 2.85
A32_M1, A32_M2, A32_M32 = 4.55, 5.325, 6.05
A33_M1, A33_M33 = 7.85, 8.85
# M1 -> A31, A32, A33
seg(A31_M1, M1y, A33_M1, M1y)
path([(A31_M1,M1y),(A31_M1,BOTF)]); mtag('M1', A31_M1, 4.54)
path([(A32_M1,M1y),(A32_M1,BOTF)]); mtag('M1', A32_M1, 4.54)
path([(A33_M1,M1y),(A33_M1,BOTF)]); mtag('M1', A33_M1, 4.54)
# singles
path([(A31_M31,SING),(A31_M31,BOTF)]); mtag('M31', A31_M31, 4.54)
path([(A32_M2,SING),(A32_M2,BOTF)]);   mtag('M2',  A32_M2, 4.54)
path([(A32_M32,SING),(A32_M32,BOTF)]); mtag('M32', A32_M32, 4.54)
path([(A33_M33,SING),(A33_M33,BOTF)]); mtag('M33', A33_M33, 4.54)
# long labels (legend)
add_label('M1 Platforms hosting ATMOS node compute',0.45,6.95,1.95,0.48,10, fill_white=True)
add_label('M31 ATMOS uncertainty module; sensor metadata',2.45,6.95,2.05,0.48,10, fill_white=True)
add_label('M2 ABLE-LBM / reduced models',4.55,6.95,1.55,0.48,10, fill_white=True)
add_label('M32 ABLE-LBM uncertainty routines',6.15,6.95,1.85,0.48,10, fill_white=True)
add_label('M33 ATMOS packaging/serialization; DDS data types',8.02,6.95,2.30,0.48,10, fill_white=True)

prs.save('ATMOS_A3_IDEF0_native_vNext_corrected.pptx')
print('saved; shapes:', len(slide.shapes._spTree))
