#!/usr/bin/env python3
"""
Generate ATMOS_A7_IDEF0_native_vNext.pptx
A7 EXTERNAL-REFERENCE IDEF0 (2x2: A71/A72 upper, A73/A74 lower).
11x8.5 landscape, reference frame style, bare 'A7' node id, external note.
>=10pt. External C2/mission-system exchanges; ATMOS not a decision authority.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK=RGBColor(0,0,0); WHITE=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(11.0); prs.slide_height=Inches(8.5)
slide=prs.slides.add_slide(prs.slide_layouts[6]); shapes=slide.shapes
def IN(v): return Inches(v)

def add_box(node,name,l,t,w,h):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=WHITE
    sp.line.color.rgb=BLACK; sp.line.width=Pt(1.5); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0=tf.paragraphs[0]; p0.alignment=PP_ALIGN.CENTER
    r0=p0.add_run(); r0.text=node; r0.font.size=Pt(12); r0.font.bold=True; r0.font.color.rgb=BLACK
    p1=tf.add_paragraph(); p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text=name; r1.font.size=Pt(10); r1.font.color.rgb=BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb=shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb=WHITE; tb.line.fill.background()
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=BLACK
    return tb

PORT=0.035
def add_port(x,y):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit=False

def _arrow(c):
    ln=c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c=shapes.add_connector(MSO_CONNECTOR.ELBOW,IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb=BLACK; c.line.width=Pt(1.25); c.shadow.inherit=False
    if arrow:_arrow(c)
    return c

def path(pts,src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code,x,y):
    add_label(code,x-0.33,y,0.66,0.24,10,bold=True,align=PP_ALIGN.CENTER,fill_white=True)

# ---------------- frame ----------------
bd=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb=BLACK; bd.line.width=Pt(1.25); bd.shadow.inherit=False
add_label('A7 — External Reference: Execute Mission Adaptation',
          0.35,0.60,10.30,0.32,15,bold=True,align=PP_ALIGN.CENTER)
add_label('EXTERNAL TO ATMOS SYSTEM BOUNDARY',0.55,0.96,4.5,0.28,11,bold=True)
nb=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(9.35),IN(7.62),IN(1.20),IN(0.30))
nb.fill.background(); nb.line.color.rgb=BLACK; nb.line.width=Pt(1.0); nb.shadow.inherit=False
nt=nb.text_frame; nt.vertical_anchor=MSO_ANCHOR.MIDDLE; nt.paragraphs[0].alignment=PP_ALIGN.CENTER
nr=nt.paragraphs[0].add_run(); nr.text='A7'; nr.font.size=Pt(12); nr.font.bold=True; nr.font.color.rgb=BLACK

# ---------------- boxes ----------------
add_box('A71','Exchange Mission Adjustment Commands',1.35,2.55,2.25,0.90)
add_box('A72','Publish Vertical Operations Execution Updates',6.40,2.55,2.35,0.90)
add_box('A73','Publish UAS Tasking Updates',1.35,5.05,2.25,0.90)
add_box('A74','Publish Coordination Status',6.40,5.05,2.35,0.90)
RBX=10.55

# ---------------- INPUTS (left) ----------------
# External Mission Decisions -> A71, A73 (left bus)
seg(0.50,4.15,0.90,4.15); seg(0.90,3.00,0.90,5.30)
path([(0.90,3.00),(1.35,3.00)])                               # -> A71 left
path([(0.90,5.30),(1.35,5.30)])                               # -> A73 left
add_label('External Mission Decisions',0.40,3.66,0.95,0.70,10,fill_white=True)
# Aircraft Execution State -> A72 left (via row gap)
path([(0.50,3.60),(6.15,3.60),(6.15,3.00),(6.40,3.00)])
add_label('Aircraft Execution State',4.55,3.18,1.70,0.32,10,fill_white=True)
# Coordination State -> A74 left (via row gap)
path([(0.50,4.35),(6.15,4.35),(6.15,5.25),(6.40,5.25)])
add_label('Coordination State',4.55,4.92,1.55,0.32,10,fill_white=True)

# ---------------- CONTROLS (top) ----------------
path([(2.475,1.50),(2.475,2.55)])                             # C71 -> A71
path([(7.575,1.50),(7.575,2.55)])                             # C72 -> A72
path([(4.05,1.50),(4.05,4.60),(2.475,4.60),(2.475,5.05)])    # C73 -> A73 top (around A71)
path([(5.95,1.50),(5.95,4.60),(7.575,4.60),(7.575,5.05)])   # C74 -> A74 top (around A72)
add_label('C71 Command authority; C2 rules of engagement',1.30,1.00,2.30,0.42,10)
add_label('C72 Reporting policy; communications availability',6.30,1.00,2.45,0.42,10)
add_label('C73 Mission authority; GCS procedures',3.10,2.05,1.90,0.42,10,fill_white=True)
add_label('C74 Coordination policy; reporting cadence',5.15,2.05,2.15,0.42,10,fill_white=True)

# ---------------- OUTPUTS (right) ----------------
path([(3.60,2.78),(3.95,2.78),(3.95,3.86),(RBX,3.86)],src_port=True)   # MAC A71->boundary
path([(8.75,2.78),(RBX,2.78)],src_port=True)                          # VOEU A72->boundary
path([(3.60,5.30),(4.20,5.30),(4.20,4.12),(RBX,4.12)],src_port=True)  # UTU A73->boundary
path([(8.75,5.50),(RBX,5.50)],src_port=True)                          # CSU A74->boundary
add_label('Mission Adjustment Commands (IER-19)',8.85,3.58,1.70,0.40,10,fill_white=True)
add_label('Vertical Ops Execution Updates (IER-20)',8.85,2.36,1.70,0.40,10,fill_white=True)
add_label('UAS Tasking Updates (IER-21)',8.95,4.20,1.60,0.34,10,fill_white=True)
add_label('Coordination Status Updates (IER-22)',8.85,5.08,1.70,0.40,10,fill_white=True)

# ---------------- INTERNAL FLOWS ----------------
path([(3.60,3.18),(3.90,3.18),(3.90,6.45),(0.95,6.45),(0.95,5.70),(1.35,5.70)],src_port=True)  # A7-F1 A71->A73
path([(3.60,5.70),(4.80,5.70),(4.80,5.50),(6.40,5.50)],src_port=True)                          # A7-F2 A73->A74
path([(8.75,3.18),(9.05,3.18),(9.05,6.45),(6.05,6.45),(6.05,5.75),(6.40,5.75)],src_port=True)  # A7-F3 A72->A74
add_label('A7-F1 Mission Decisions',1.55,6.50,1.55,0.30,10,fill_white=True)
add_label('A7-F2 Tasking Updates',4.85,5.74,1.45,0.30,10,fill_white=True)
add_label('A7-F3 Execution Updates',6.20,6.50,1.55,0.30,10,fill_white=True)

# ---------------- MECHANISM-LANE MATRIX ----------------
BUS=6.80
# M71 -> A73 (direct) and A71 (around lower row, via central riser)
path([(2.30,BUS),(2.30,5.95)]); mtag('M71',2.30,6.02)
path([(4.45,BUS),(4.45,4.85),(2.55,4.85),(2.55,3.45)]); mtag('M71',2.55,3.55)
# M72 -> A72 (around lower row)
path([(8.85,BUS),(8.85,4.85),(7.10,4.85),(7.10,3.45)]); mtag('M72',7.10,3.55)
# M73 -> A74 (direct) and A72 (around lower row)
path([(7.55,BUS),(7.55,5.95)]); mtag('M73',7.55,6.02)
path([(9.15,BUS),(9.15,4.70),(7.80,4.70),(7.80,3.45)]); mtag('M73',7.80,3.55)
add_label('M71 C2 systems / TOC',1.20,6.95,2.0,0.30,10,fill_white=True)
add_label('M72 Aircraft avionics',4.30,6.95,1.9,0.30,10,fill_white=True)
add_label('M73 Communications network / DDS',6.95,6.95,2.3,0.30,10,fill_white=True)

prs.save('ATMOS_A7_IDEF0_native_vNext.pptx')
print('saved; shapes:',len(slide.shapes._spTree))
