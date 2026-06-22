#!/usr/bin/env python3
"""Build ATMOS_A0_native_clean.pptx — native, editable IDEF0 A0 decomposition.

Coordinate-placed right-angle connector paths (arrowhead on the terminal segment,
landing exactly on each box's outside-perimeter anchor). No glue -> no risk of an
arrowhead snapping inside a box. Native rectangles + text boxes + connectors only;
no images, no groups spanning the whole diagram.

Model source: ATMOS_A0.yaml. Internal flows F1-F5 only. Black-and-white IDEF0.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(12.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.34, 0.45, 0.28, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

# -------- frame --------
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
                            Inches(21.0), Inches(10.7))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# -------- title / classification / node --------
label(7.0, 0.22, 10.0, 0.45, "A0 — Decompose Conduct Air-Focused Weather Exploitation Operations",
      size=18, bold=True, align=PP_ALIGN.CENTER)
label(17.6, 0.26, 3.8, 0.3, "Distribution Statement C", size=12, align=PP_ALIGN.RIGHT)
label(19.5, 11.42, 1.9, 0.3, "Node: A0", size=12, bold=True, align=PP_ALIGN.RIGHT)

# -------- activity boxes --------
box(2.4, 5.2, 2.6, 1.3, "Acquire Micro-Weather Observations", "A1")
box(5.7, 5.2, 2.6, 1.3, "Generate Local Weather State", "A2")
box(9.0, 5.2, 2.6, 1.3, "Quantify Uncertainty", "A3")
box(12.3, 5.2, 2.8, 1.3, "Federate and Maintain Federated Weather Context", "A4")
box(16.4, 3.0, 3.0, 1.3, "Detect Threshold Crossings (Descriptive Support)", "A5")
box(16.4, 7.4, 3.0, 1.3, "Produce Mission-Tailored Weather Context", "A6")

# -------- internal flows F1-F5 --------
path([(5.0, 5.85), (5.7, 5.85)])
path([(8.3, 5.85), (9.0, 5.85)])
path([(11.6, 5.85), (12.3, 5.85)])
path([(15.1, 5.55), (15.7, 5.55), (15.7, 3.65), (16.4, 3.65)])    # F4 -> A5 left
path([(15.1, 6.15), (15.7, 6.15), (15.7, 7.85), (16.4, 7.85)])    # F5 -> A6 left
label(4.05, 4.55, 2.6, 0.5, "F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)", size=9, align=PP_ALIGN.CENTER)
label(7.35, 4.35, 2.6, 0.5, "F2 Local Micro-Weather State Estimate (IER-05)", size=9, align=PP_ALIGN.CENTER)
label(10.6, 4.55, 2.6, 0.5, "F3 Confidence Bounds & Risk Envelopes (IER-09)", size=9, align=PP_ALIGN.CENTER)
label(14.9, 2.45, 2.3, 0.5, "F4 Federated Weather Context / COWP State", size=9)
label(10.7, 7.12, 2.6, 0.5, "F5 Federated Weather Context / COWP State", size=9)

# -------- inputs (LEFT boundary) --------
path([(0.5, 5.55), (2.4, 5.55)])           # I1 -> A1
path([(0.5, 6.15), (2.4, 6.15)])           # I2 -> A1
path([(0.5, 9.8), (15.4, 9.8)], arrow=False)            # I3 bus
path([(15.4, 9.8), (15.4, 4.05), (16.4, 4.05)])         # I3 -> A5
path([(15.4, 9.8), (15.4, 8.25), (16.4, 8.25)])         # I3 -> A6
label(0.6, 5.18, 1.75, 0.34, "I1 Collocated Atmospheric Observations", size=9)
label(0.6, 6.18, 1.75, 0.34, "I2 Peer Platform Weather Observations", size=9)
label(0.6, 9.45, 2.9, 0.34, "I3 Mission Weather Threshold Definitions", size=9)

# -------- controls (TOP boundary) --------
path([(13.3, 1.1), (13.3, 5.2)])                                       # C3 -> A4
path([(14.2, 1.1), (14.2, 1.4), (20.0, 1.4)], arrow=False)             # C2 trunk
path([(14.2, 1.4), (14.2, 5.2)])                                       # C2 -> A4
path([(20.0, 1.4), (20.0, 6.5), (17.6, 6.5), (17.6, 7.4)])            # C2 -> A6
path([(17.9, 1.1), (17.9, 1.5), (20.6, 1.5)], arrow=False)            # C1 trunk
path([(17.9, 1.5), (17.9, 3.0)])                                       # C1 -> A5
path([(20.6, 1.5), (20.6, 6.9), (18.3, 6.9), (18.3, 7.4)])           # C1 -> A6
label(10.5, 0.66, 2.6, 0.34, "C3 Network availability and DDS QoS policies", size=9, align=PP_ALIGN.CENTER)
label(13.6, 0.66, 2.5, 0.34, "C2 OPSEC / classification guidance", size=9, align=PP_ALIGN.CENTER)
label(16.9, 0.66, 3.1, 0.34, "C1 Mission objectives and operational constraints", size=9, align=PP_ALIGN.CENTER)

# -------- outputs (to RIGHT boundary) --------
path([(15.1, 5.85), (21.5, 5.85)])         # O1 from A4 (middle lane)
path([(19.4, 3.65), (21.5, 3.65)])         # O3 from A5
path([(19.4, 8.05), (21.5, 8.05)])         # O2 from A6
label(21.6, 5.60, 2.3, 0.5, "O1 Fused COWP State", size=9)
label(21.6, 3.40, 2.3, 0.5, "O3 Weather State Change Notifications", size=9)
label(21.6, 7.80, 2.3, 0.5, "O2 Mission-Tailored COWP Excerpts", size=9)

# -------- mechanisms (BOTTOM boundary) --------
# M1 -> A1..A6
path([(4.0, 11.8), (4.0, 10.5)], arrow=False)           # M1 stub
path([(3.7, 10.5), (17.9, 10.5)], arrow=False)          # M1 bus
path([(3.7, 10.5), (3.7, 6.5)])                         # M1 -> A1
path([(6.6, 10.5), (6.6, 6.5)])                         # M1 -> A2
path([(9.9, 10.5), (9.9, 6.5)])                         # M1 -> A3
path([(13.5, 10.5), (13.5, 6.5)])                       # M1 -> A4
path([(17.6, 10.5), (17.6, 8.7)])                       # M1 -> A6
path([(15.9, 10.5), (15.9, 5.9), (17.9, 5.9), (17.9, 4.3)])   # M1 -> A5 (around A6 via gap)
# M2 -> A2, A3
path([(7.0, 11.8), (7.0, 10.9)], arrow=False)           # M2 stub
path([(7.3, 10.9), (10.6, 10.9)], arrow=False)          # M2 bus
path([(7.3, 10.9), (7.3, 6.5)])                         # M2 -> A2
path([(10.6, 10.9), (10.6, 6.5)])                       # M2 -> A3
# M3 -> A4, A6
path([(14.0, 11.8), (14.0, 11.1)], arrow=False)         # M3 stub
path([(14.3, 11.1), (18.3, 11.1)], arrow=False)         # M3 bus
path([(14.3, 11.1), (14.3, 6.5)])                       # M3 -> A4
path([(18.3, 11.1), (18.3, 8.7)])                       # M3 -> A6
label(2.0, 11.9, 3.2, 0.32, "M1 Platforms hosting ATMOS node compute", size=9, align=PP_ALIGN.CENTER)
label(5.8, 11.9, 3.0, 0.32, "M2 ABLE-LBM / reduced models", size=9, align=PP_ALIGN.CENTER)
label(12.4, 11.9, 3.4, 0.32, "M3 DDS middleware and communications links", size=9, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A0_native_clean.pptx")
print("saved ATMOS_A0_native_clean.pptx; shapes:", len(slide.shapes))
