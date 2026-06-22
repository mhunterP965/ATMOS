#!/usr/bin/env python3
"""
Generate ATMOS_A5_IDEF0_native_vNext.pptx
A5 IDEF0 decomposition: chain A51->A52->A53->A54. 11x8.5 landscape, reference
frame style, bare 'A5' node id. F4 input bus -> A52/A53/A54 left; I3 -> A51 left;
A5-F1 internal control A51 right -> A52 top. Auditable mechanism matrix. >=10pt.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK=RGBColor(0,0,0); WHITE=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(11.0); prs.slide_height=Inches(8.5)
slide=prs.slides.add_slide(prs.slide_layouts[6]); shapes=slide.shapes
def IN(v): return Inches(v)

def add_box(node,name,l,t,w,h):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=WHITE
    sp.line.color.rgb=BLACK; sp.line.width=Pt(1.5); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in('margin_left','margin_right'): setattr(tf,m,Pt(2))
    for m in('margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p0=tf.paragraphs[0]; p0.alignment=PP_ALIGN.CENTER
    r0=p0.add_run(); r0.text=node; r0.font.size=Pt(12); r0.font.bold=True; r0.font.color.rgb=BLACK
    p1=tf.add_paragraph(); p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text=name; r1.font.size=Pt(10); r1.font.color.rgb=BLACK
    return sp

def add_label(text,l,t,w,h,size=10,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,fill_white=False):
    tb=shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    if fill_white:
        tb.fill.solid(); tb.fill.fore_color.rgb=WHITE; tb.line.fill.background()
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,Pt(1))
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=BLACK
    return tb

PORT=0.035
def add_port(x,y):
    sp=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(x-PORT/2),IN(y-PORT/2),IN(PORT),IN(PORT))
    sp.fill.background(); sp.line.fill.background(); sp.shadow.inherit=False

def _arrow(c):
    ln=c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))

def seg(x1,y1,x2,y2,arrow=False):
    c=shapes.add_connector(MSO_CONNECTOR.ELBOW,IN(x1),IN(y1),IN(x2),IN(y2))
    c.line.color.rgb=BLACK; c.line.width=Pt(1.25); c.shadow.inherit=False
    if arrow:_arrow(c)
    return c

def path(pts,src_port=False):
    for i in range(len(pts)-1):
        seg(*pts[i],*pts[i+1],arrow=(i==len(pts)-2))
    add_port(*pts[-1])
    if src_port: add_port(*pts[0])

def mtag(code,x,y):
    add_label(code,x-0.33,y,0.66,0.24,10,bold=True,align=PP_ALIGN.CENTER,fill_white=True)

# ---------------- frame ----------------
bd=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(0.35),IN(0.55),IN(10.30),IN(7.45))
bd.fill.background(); bd.line.color.rgb=BLACK; bd.line.width=Pt(1.25); bd.shadow.inherit=False
add_label('A5 — Decompose Detect Threshold Crossings (Descriptive Support)',
          0.35,0.60,10.30,0.34,15,bold=True,align=PP_ALIGN.CENTER)
nb=shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(9.35),IN(7.62),IN(1.20),IN(0.30))
nb.fill.background(); nb.line.color.rgb=BLACK; nb.line.width=Pt(1.0); nb.shadow.inherit=False
nt=nb.text_frame; nt.vertical_anchor=MSO_ANCHOR.MIDDLE; nt.paragraphs[0].alignment=PP_ALIGN.CENTER
nr=nt.paragraphs[0].add_run(); nr.text='A5'; nr.font.size=Pt(12); nr.font.bold=True; nr.font.color.rgb=BLACK

# ---------------- boxes ----------------
add_box('A51','Receive Mission Weather Threshold Definitions',1.05,3.45,2.00,0.90)
add_box('A52','Present Weather Conditions Relevant to Aviation Risk',3.35,3.45,2.05,0.90)
add_box('A53','Present Temporal Weather Condition Trends',5.70,3.45,2.05,0.90)
add_box('A54','Notify Weather State Changes',8.05,3.45,2.00,0.90)
CY=3.90; TOPF=3.45; BOTF=4.35; RBX=10.45

# ---------------- INPUTS (left) ----------------
path([(0.50,CY),(1.05,CY)])                                   # I3 -> A51 left
add_label('I3 Mission Weather Threshold Definitions (IER-14)',0.40,3.02,1.60,0.42,10,fill_white=True)
# F4 input bus -> A52,A53,A54 (separate left ports @3.65)
seg(0.50,2.60,7.90,2.60)
path([(3.28,2.60),(3.28,3.65),(3.35,3.65)])                   # F4 -> A52 left
path([(5.55,2.60),(5.55,3.65),(5.70,3.65)])                   # F4 -> A53 left
path([(7.90,2.60),(7.90,3.65),(8.05,3.65)])                   # F4 -> A54 left
add_label('F4 Federated Weather Context / COWP State',0.40,2.22,1.95,0.36,10,fill_white=True)

# ---------------- CONTROLS (top) ----------------
path([(2.05,1.55),(2.05,TOPF)])                               # C51 -> A51
seg(3.70,1.45,6.10,1.45)                                      # C1 bus
path([(3.70,1.45),(3.70,TOPF)])                               # C1 -> A52
path([(6.10,1.45),(6.10,TOPF)])                               # C1 -> A53
path([(4.30,1.62),(4.30,TOPF)])                               # C52 -> A52
path([(6.85,1.62),(6.85,TOPF)])                               # C53 -> A53
path([(8.55,1.55),(8.55,TOPF)])                               # C2 -> A54
path([(9.35,1.62),(9.35,TOPF)])                               # C54 -> A54
add_label('C51 Mission configuration management',0.95,1.00,1.95,0.42,10)
add_label('C1 Mission objectives and operational constraints',3.00,0.98,2.15,0.42,10)
add_label('C52 Presentation rules; role/platform relevance filters',2.95,1.40,2.45,0.40,10)
add_label('C53 Trend window definitions; smoothing/aggregation rules',5.45,1.40,2.55,0.40,10)
add_label('C2 OPSEC / classification guidance',7.40,1.00,1.60,0.42,10)
add_label('C54 Notification thresholds; communications/QoS policy',9.05,0.98,1.55,0.55,10)

# ---------------- INTERNAL FLOWS ----------------
# A5-F1 internal CONTROL: A51 right -> A52 TOP (vertical perpendicular)
path([(3.05,CY),(3.12,CY),(3.12,2.95),(4.90,2.95),(4.90,TOPF)],src_port=True)
add_label('A5-F1 Threshold Parameters Available for Comparison',3.30,2.55,2.05,0.36,10,fill_white=True)
# A5-F2 A52 right -> A53 left ; A5-F3 A53 right -> A54 left
path([(5.40,CY),(5.70,CY)],src_port=True)
path([(7.75,CY),(8.05,CY)],src_port=True)
add_label('A5-F2 Weather Conditions Relevant to Aviation Risk (IER-15)',4.55,4.52,2.00,0.42,10,fill_white=True)
add_label('A5-F3 Temporal Weather Condition Trends (IER-16)',6.95,4.52,2.00,0.42,10,fill_white=True)

# ---------------- OUTPUT (right) ----------------
path([(10.05,CY),(RBX,CY)],src_port=True)                     # O3 -> right boundary
add_label('O3 Weather State Change Notifications (IER-17)',8.55,2.92,2.05,0.40,10,fill_white=True)
add_label('To Platforms / TOC / Downstream A6',8.85,4.50,1.75,0.40,10,fill_white=True)

# ---------------- MECHANISM-LANE MATRIX ----------------
SING=5.25; M3y=5.65; M1y=6.05
# bottom ports
A51_M1,A51_M51=1.55,2.55
A52_M1,A52_M3,A52_M52=3.75,4.35,4.95
A53_M1,A53_M53=6.10,7.10
A54_M1,A54_M3,A54_M54=8.55,9.10,9.65
# M1 spine (all four)
seg(A51_M1,M1y,A54_M1,M1y)
for x in (A51_M1,A52_M1,A53_M1,A54_M1):
    path([(x,M1y),(x,BOTF)]); mtag('M1',x,4.44)
# M3 spine (A52,A54)
seg(A52_M3,M3y,A54_M3,M3y)
path([(A52_M3,M3y),(A52_M3,BOTF)]); mtag('M3',A52_M3,4.44)
path([(A54_M3,M3y),(A54_M3,BOTF)]); mtag('M3',A54_M3,4.44)
# singles
path([(A51_M51,SING),(A51_M51,BOTF)]); mtag('M51',A51_M51,4.44)
path([(A52_M52,SING),(A52_M52,BOTF)]); mtag('M52',A52_M52,4.44)
path([(A53_M53,SING),(A53_M53,BOTF)]); mtag('M53',A53_M53,4.44)
path([(A54_M54,SING),(A54_M54,BOTF)]); mtag('M54',A54_M54,4.44)
# long labels (legend)
add_label('M1 Platforms hosting ATMOS node compute',0.45,6.40,1.70,0.55,10,fill_white=True)
add_label('M3 DDS middleware and communications links',2.20,6.40,1.85,0.55,10,fill_white=True)
add_label('M51 TOC/GCS/platform mission configuration interface; ATMOS parameter store',4.10,6.40,1.75,0.95,10,fill_white=True)
add_label('M52 ATMOS dissemination/presentation logic',5.95,6.40,1.70,0.75,10,fill_white=True)
add_label('M53 ATMOS temporal analytics; local compute',7.70,6.40,1.65,0.75,10,fill_white=True)
add_label('M54 DDS events/notifications; communications subsystem',9.40,6.40,1.20,0.95,10,fill_white=True)

prs.save('ATMOS_A5_IDEF0_native_vNext.pptx')
print('saved; shapes:',len(slide.shapes._spTree))
