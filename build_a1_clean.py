#!/usr/bin/env python3
"""Build ATMOS_A1_native_clean.pptx — native, editable IDEF0 A1 decomposition.

Model (authoritative, inline from the task): parent A1 'Acquire Micro-Weather
Observations' -> A11, A12, A13. Internal flows A1-F1, A1-F2 only. A14/TBD NOT drawn.
Coordinate-placed right-angle connector paths (arrowhead on the terminal segment,
exactly on each box's outside-perimeter anchor). No glue, no images, no whole-diagram group.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(24)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def seg(x1, y1, x2, y2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5)
    if arrow:
        ln = c.line._get_or_add_ln()
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def path(points, arrow=True):
    for i in range(len(points) - 1):
        x1, y1 = points[i]; x2, y2 = points[i + 1]
        seg(x1, y1, x2, y2, arrow=(arrow and i == len(points) - 2))


def label(x, y, w, h, text, size=10, align=PP_ALIGN.LEFT, bold=False,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK; r.font.name = "Arial"
    return tb


def box(x, y, w, h, name, desc, nid):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BLACK; rect.line.width = Pt(2.0); rect.shadow.inherit = False
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(9); r2.font.color.rgb = BLACK; r2.font.name = "Arial"
    label(x + w - 0.55, y + h - 0.32, 0.45, 0.26, nid, size=11, bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    return rect

# -------- frame --------
fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
                            Inches(21.0), Inches(10.7))
fr.fill.background(); fr.line.color.rgb = BLACK; fr.line.width = Pt(1.25); fr.shadow.inherit = False

# -------- title / marking / node --------
label(6.0, 0.22, 11.0, 0.45, "A1 — Decompose Acquire Micro-Weather Observations",
      size=18, bold=True, align=PP_ALIGN.CENTER)
label(17.6, 0.26, 3.8, 0.3, "Distribution Statement C", size=12, align=PP_ALIGN.RIGHT)
label(19.5, 11.42, 1.9, 0.3, "Node: A1", size=12, bold=True, align=PP_ALIGN.RIGHT)

# -------- activity boxes (A11 upper-left, A12 lower-left offset, A13 center-right) --------
box(2.6, 3.0, 3.0, 1.4, "Sense Collocated Atmospheric Conditions",
    "Collect local atmospheric observations using onboard or attached sensors.", "A11")
box(3.8, 7.4, 3.0, 1.4, "Receive Peer Platform Weather Observations",
    "Ingest weather observations from peer air/ground nodes via DDS when connectivity permits.", "A12")
box(12.5, 5.2, 3.0, 1.4, "Time-Tag, Geo-Tag, and Quality-Check Observations",
    "Apply time, geolocation, and quality control processing to ensure consistency and traceability.", "A13")
# A11: 2.6-5.6, 3.0-4.4 | A12: 3.8-6.8, 7.4-8.8 | A13: 12.5-15.5, 5.2-6.6

# -------- internal flows A1-F1, A1-F2 --------
path([(5.6, 3.7), (8.3, 3.7), (8.3, 5.6), (12.5, 5.6)])   # A1-F1: A11 -> A13 left
path([(6.8, 8.1), (8.7, 8.1), (8.7, 6.2), (12.5, 6.2)])   # A1-F2: A12 -> A13 left
label(6.4, 3.36, 2.4, 0.4, "A1-F1 Collocated Atmospheric Observations", size=9)
label(6.9, 8.16, 2.4, 0.4, "A1-F2 Peer Platform Weather Observations", size=9)

# -------- inputs (LEFT boundary) --------
path([(0.5, 3.7), (2.6, 3.7)])   # I1 -> A11 left
path([(0.5, 8.1), (3.8, 8.1)])   # I2 -> A12 left
label(0.6, 3.36, 2.0, 0.32, "I1 Collocated Atmospheric Observations", size=9)
label(0.6, 7.76, 2.0, 0.32, "I2 Peer Platform Weather Observations", size=9)

# -------- controls (TOP boundary) --------
path([(4.1, 1.1), (4.1, 3.0)])     # C11 -> A11 top
path([(6.2, 1.1), (6.2, 7.4)])     # C12 -> A12 top (right of A11, clear)
path([(14.0, 1.1), (14.0, 5.2)])   # C13 -> A13 top
label(1.5, 0.62, 3.0, 0.42, "C11 Sensor configuration; platform operating state; sampling policies", size=9, align=PP_ALIGN.CENTER)
label(6.0, 0.62, 3.0, 0.42, "C12 DDS subscription filters; QoS; communications availability", size=9, align=PP_ALIGN.CENTER)
label(11.0, 0.62, 3.0, 0.42, "C13 QC rules; time/position accuracy requirements", size=9, align=PP_ALIGN.CENTER)

# -------- output (parent F1 to RIGHT boundary) --------
path([(15.5, 5.9), (21.5, 5.9)])   # F1 from A13 right
label(21.6, 5.45, 2.4, 0.55, "F1 Time/Geo-Tagged, Quality-Checked Observations (IER-03)", size=9)
label(21.6, 6.15, 2.4, 0.4, "To A2 Generate Local Weather State", size=9, bold=True)

# -------- mechanisms (BOTTOM boundary) --------
# M1 -> A11, A12, A13 (shared mechanism: stub + bus + risers)
path([(6.0, 11.8), (6.0, 10.0)], arrow=False)            # M1 stub
path([(3.0, 10.0), (13.5, 10.0)], arrow=False)           # M1 bus
path([(3.0, 10.0), (3.0, 4.4)])                          # M1 -> A11 (left of A12)
path([(4.3, 10.0), (4.3, 8.8)])                          # M1 -> A12
path([(13.5, 10.0), (13.5, 6.6)])                        # M1 -> A13
# M11, M12, M13 (local mechanisms)
path([(3.3, 11.8), (3.3, 4.4)])                          # M11 -> A11
path([(5.0, 11.8), (5.0, 8.8)])                          # M12 -> A12
path([(14.5, 11.8), (14.5, 6.6)])                        # M13 -> A13
label(1.3, 11.9, 2.6, 0.32, "M11 Onboard/attached sensors", size=9, align=PP_ALIGN.CENTER)
label(4.2, 11.9, 2.8, 0.42, "M12 DDS middleware; radios/network transport", size=9, align=PP_ALIGN.CENTER)
label(7.3, 11.9, 2.9, 0.32, "M1 Platforms hosting ATMOS node compute", size=9, align=PP_ALIGN.CENTER)
label(11.6, 11.9, 3.6, 0.42, "M13 ATMOS edge processing; navigation solution (GPS/INS); local compute", size=9, align=PP_ALIGN.CENTER)

prs.save("ATMOS_A1_native_clean.pptx")
print("saved ATMOS_A1_native_clean.pptx; shapes:", len(slide.shapes))
